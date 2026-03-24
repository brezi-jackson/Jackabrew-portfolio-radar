from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

out_path = 'projects/brezi-coffee/Brezi_Investor_Talent_Deck.pptx'
prs = Presentation()
prs.slide_height = Inches(9)
prs.slide_width = Inches(16)

# Helpers
def set_bg(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    from pptx.dml.color import RGBColor
    fill.fore_color.rgb = RGBColor(*rgb)

def add_title_slide(title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, (0,0,0))
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(2.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(14), Inches(1))
        st = sub.text_frame
        st.text = subtitle
        st.paragraphs[0].font.size = Pt(20)
        st.paragraphs[0].font.color.rgb = RGBColor(200,200,200)

def add_bullets_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, (0,0,0))
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(14.8), Inches(1))
    t = title_box.text_frame
    t.text = title
    t.paragraphs[0].font.size = Pt(28)
    t.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    left = Inches(0.8)
    top = Inches(1.6)
    width = Inches(14)
    height = Inches(6)
    body = slide.shapes.add_textbox(left, top, width, height)
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i==0:
            p = tf.paragraphs[0]
            p.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(220,220,220)

# Build slides
add_title_slide('Brezi — Start Cold.', 'Brezi: The Cold Protocol Platform\nInvestor & Talent Deck')

add_bullets_slide('The Problem', [
    'Morning ritual is the last unoptimized input — people optimize sleep, training, nutrition but not the first thing they consume.',
    'Coffee consumption is variable, acidic, and often disrupts recovery metrics (HRV/cortisol).',
    'No brand connects cold exposure practices to morning consumption in a meaningful way.'
])

add_bullets_slide('Our Core Insight', [
    'Cold is a deliberate physiological practice — not just a temperature preference.',
    'Brezi connects the morning consumption ritual to cold wellness and biometric data.',
    'Coffee is the entry product; cold is the platform.'
])

add_bullets_slide('Solution — Brezi', [
    'PrecisionBrew™: a thermoelectric cold-extraction machine (12-minute brew ritual).',
    'Brezi OS: device + Brezi Health Graph + BreziIntel (personalized recommendations).',
    'Platform: hardware + consumables + app integrations + community + spaces.'
])

add_bullets_slide('Signature Ritual — The 12-Minute Window', [
    'The brew takes 12 minutes — that is the product and ritual.',
    'A single moving cold-blue indicator marks the 12-minute cycle — the owner’s time.',
    'All brand moments and extensions organize around this window (content, kits, spaces).'
])

add_bullets_slide('Target Customer — Protocol Optimizer', [
    'Aged 27–38, uses Oura/WHOOP, serious about morning routines.',
    'Performs Zone 2 / functional fitness, follows Huberman/Peter Attia.',
    'Buys identity & optimized results, not commodity appliances.'
])

add_bullets_slide('Brand & Aesthetic', [
    'Precision engineering × monastic calm: matte materials, glacier blue accent.',
    'Counterpoint of bone/ivory white, photography in pre-dawn light.',
    'Object-first design: the countertop test — an object people display.'
])

add_bullets_slide('Tech Vision — Brezi OS', [
    'Device Layer: telemetry, OTA, cartridge NFC auth.',
    'Data Layer: Brezi Health Graph — HRV, sleep, brew sessions, subjective feedback.',
    'Insights Layer: BreziIntel — personalized brew recommendations & longitudinal correlations.'
])

add_bullets_slide('Moat & Differentiation', [
    'Unique dataset: biometric state at brew moment + brew parameters + 2h performance.',
    'Cartridge NFC + formulation IDs = product + data lock-in.',
    'Local-first ML (privacy-first) + patentable thermal profiling algorithms and university research.'
])

add_bullets_slide('Product Roadmap', [
    'Now: Ship PrecisionBrew™ + Apple Health integration + manual HRV input.',
    '18 months: BreziRest Pod — evening recovery cold-extract system (consumable model).',
    '30 months: BreziBar — cold-extraction station for adaptogens & functional compounds.'
])

add_bullets_slide('Go-to-Market — The Cold Protocol', [
    'Launch film & campaign: "Start Cold" hero film (60–90s).',
    '30-day Cold Protocol Challenge — app + UGC + WHOOP/Oura correlations.',
    'Invite-only Cold Morning pop-up in HK (no sale; brand signaling & press).' 
])

add_bullets_slide('Ecosystem & Revenue', [
    'Hardware retail (~$350–400) + consumable subscription (cartridges).',
    'B2B: premium gyms, hotels, corporate protocol offering.',
    'Long term: data and platform services, partnerships with wearables/health apps.'
])

add_bullets_slide('Traction & Assets', [
    'Kickstarter launch success; Shopify store live; Taobao presence.',
    'Completed Phase 1 research: 7 interviews + 181 backer surveys (data available).',
    'iOS technical spec ready; UX copy brief done; Jensen coding pending design lock.'
])

add_bullets_slide('Team', [
    'Founders: Jackson (UC Berkeley) + Graves Ho — founder DNA + product vision.',
    'Engineering: semiconductor + consumer electronics background (precision skills).',
    'Task Force: CMO, CTO, Creative (Bobo), PM (coordinator) — assembled to build the brand.'
])

add_bullets_slide('The Ask — Talent & Capital', [
    'Talent: Senior hardware engineer (thermal systems), ML engineer (privacy-first models), product designer (industrial + UX).',
    'Capital: $X for 18 months (complete precision batch, app integrations, BreziRest R&D, launch campaign).',
    'Partnerships: Oura/WHOOP integrations, research partners (HKUST/HKU), creative collabs (Satisfy Running, Aesop/Byredo).'
])

add_bullets_slide('Next Steps', [
    'Finalize visual assets (product photography, hero film treatment).',
    'Design lock for iOS and immediate Jensen sprint.',
    'Prepare investor one-pager and financial model (I can draft next).'
])

add_bullets_slide('Contact & Follow-up', [
    'Jackson Hung — Founder',
    'Brezi — projects/brezi-coffee/Brezi_Investor_Talent_Deck.pptx',
    'Request: Review & feedback. I will iterate with the task force and produce final investor-ready deck.'
])

# Save
prs.save(out_path)
print('WROTE', out_path)

"""
Brezi VC Pitch Deck v3 — Full Platform Vision
"""
import os, io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyArrowPatch
import matplotlib.ticker as mticker
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ───────────────────────────────────────────────────────────────────
BLACK = RGBColor(0,   0,   0)
WHITE = RGBColor(255, 255, 255)
GREY  = RGBColor(245, 245, 245)
BONE  = RGBColor(240, 235, 225)
TEAL  = RGBColor(43,  189, 186)
CYAN2 = RGBColor(26,  138, 136)
MONO  = RGBColor(140, 140, 140)

HEX = dict(
    black='#000000', white='#FFFFFF', grey='#F5F5F5',
    bone='#F0EBE1',  teal='#2BBDBA', cyan2='#1A8A88',
    mono='#8C8C8C',  dark='#111111', mid='#222222',
)

OUT = 'projects/brezi-coffee/Brezi_VC_Pitch_Deck_v3.pptx'
os.makedirs('projects/brezi-coffee', exist_ok=True)

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def bg(slide, rgb):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = rgb

def tb(slide, text, x, y, w, h, size=17, bold=False,
       color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box

def eyebrow(slide, text, y=0.25):
    tb(slide, text, 0.6, y, 14, 0.4, size=10, color=TEAL)

def h1(slide, text, y=0.65, size=44, color=BLACK):
    tb(slide, text, 0.6, y, 14.8, 2.0, size=size, bold=True, color=color)

def bar(slide, x, y, h, w=0.05):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = TEAL; s.line.fill.background()

def secnum(slide, n):
    tb(slide, n, 0.6, 0.3, 2, 0.5, size=11, color=TEAL)

def embed(slide, fig, x, y, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                transparent=True)
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))
    plt.close(fig)

# ════════════════════════════════════════════════════════════════════════════
# CHARTS
# ════════════════════════════════════════════════════════════════════════════

def ch_tam():
    fig, ax = plt.subplots(figsize=(8, 3))
    fig.patch.set_facecolor(HEX['grey']); ax.set_facecolor(HEX['grey'])
    labels = ['TAM — Cold Wellness +\nWearables + Premium Coffee',
              'SAM — Protocol Optimizer\nAddrессable Market',
              'SOM — 3–5 Year\nCapture']
    vals = [89.3, 0.694, 0.0135]
    colors = [HEX['cyan2'], HEX['teal'], '#48DDDA']
    bars = ax.barh(labels, vals, color=colors, height=0.5)
    for b, v in zip(bars, vals):
        lbl = f'${v:.1f}B' if v >= 1 else f'${v*1000:.0f}M'
        ax.text(b.get_width()+0.5, b.get_y()+b.get_height()/2, lbl,
                va='center', color=HEX['black'], fontsize=11, fontweight='bold')
    ax.set_xlim(0, 107); ax.spines[:].set_visible(False)
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_: f'${x:.0f}B'))
    ax.tick_params(colors=HEX['black'], labelsize=9)
    plt.tight_layout(); return fig

def ch_why_now():
    fig, ax = plt.subplots(figsize=(7.5, 3.2))
    fig.patch.set_facecolor(HEX['black']); ax.set_facecolor(HEX['black'])
    yr = [2020,2021,2022,2023,2024,2025,2026]
    ax.plot(yr,[100,118,145,178,210,255,310],color=HEX['teal'], lw=2.5,marker='o',ms=5,label='Cold Wellness')
    ax.plot(yr,[100,110,128,148,165,190,218],color='#AAAAAA',  lw=2.5,marker='s',ms=5,label='Wearables (volume)')
    ax.plot(yr,[100,112,130,155,185,220,262],color=HEX['white'],lw=2.5,marker='^',ms=5,label='Biohacking market')
    ax.axvline(2025, color=HEX['teal'], lw=1.2, linestyle='--', alpha=0.7)
    ax.text(2025.1, 275, 'Brezi\nenters', color=HEX['teal'], fontsize=8)
    ax.set_ylabel('Indexed Growth (2020=100)', color='#AAAAAA', fontsize=8)
    ax.tick_params(colors='#AAAAAA', labelsize=8)
    ax.legend(fontsize=8, facecolor='#181818', labelcolor='white', framealpha=0.9)
    for sp in ax.spines.values(): sp.set_color('#333333')
    ax.annotate('*Note: Wearables market is $87B absolute vs cold wellness $330M.\n Index shows growth rate, not size.', 
                (0,0), xycoords='axes fraction', fontsize=6.5, color='#777777',
                xytext=(0, -0.22), textcoords='axes fraction')
    plt.tight_layout(); return fig

def ch_arc():
    """Emotional arc of the day — horizontal timeline"""
    fig, ax = plt.subplots(figsize=(14, 4.5))
    fig.patch.set_facecolor(HEX['black']); ax.set_facecolor(HEX['black'])
    
    events = [
        (5.5,  'Brezi\nDeep',    'Sleep\noptimized',    '#444444'),
        (6.5,  'Brezi\nWake',    'Cold brew\nto HRV',   HEX['teal']),
        (7.0,  'Brezi\nLumen',   'Circadian\nlight',    HEX['cyan2']),
        (7.5,  'Brezi\nFuel',    'Smart\nblend',        HEX['teal']),
        (7.75, 'Brezi\nAdapt',   "Lion's mane\nextract",HEX['cyan2']),
        (9.0,  '— quiet —',      'Flow state\nno interrupt','#333333'),
        (12.5, 'Brezi\nFuel',    'Midday\nprotocol',    HEX['teal']),
        (16.0, 'Brezi\nAdapt',   'Pre-train\nrhodiola', HEX['cyan2']),
        (19.0, 'Brezi\nRestore', 'Recovery\nstack',     HEX['teal']),
        (21.0, 'Brezi\nLumen',   'Amber\nwind-down',    HEX['cyan2']),
        (22.0, 'Brezi\nRestore', 'Night\nformula',      HEX['teal']),
    ]
    
    # Timeline bar
    ax.axhline(2.0, color='#333333', lw=2, xmin=0.02, xmax=0.98)
    
    for t, prod, desc, c in events:
        quiet = prod == '— quiet —'
        ax.scatter(t, 2.0, s=120 if not quiet else 40,
                   color=c, zorder=5)
        if not quiet:
            yt = 2.9 if events.index((t,prod,desc,c)) % 2 == 0 else 1.0
            ax.plot([t, t], [2.0, yt], color=c, lw=1, alpha=0.6)
            ax.text(t, yt + (0.15 if yt > 2 else -0.3), prod,
                    ha='center', va='bottom' if yt>2 else 'top',
                    fontsize=8, color=HEX['white'], fontweight='bold')
            ax.text(t, yt + (0.55 if yt > 2 else -0.75), desc,
                    ha='center', va='bottom' if yt>2 else 'top',
                    fontsize=7, color='#AAAAAA')
        else:
            ax.text(t, 2.25, prod, ha='center', fontsize=7, color='#555555', style='italic')

    # Time labels
    for hour in [6, 9, 12, 15, 18, 21]:
        ax.text(hour, 1.55, f'{hour:02d}:00', ha='center', fontsize=8, color='#555555')
    
    ax.set_xlim(5, 23.5); ax.set_ylim(0.3, 3.8)
    ax.axis('off')
    plt.tight_layout(); return fig

def ch_ecosystem():
    """Hub-and-spoke product ecosystem"""
    fig, ax = plt.subplots(figsize=(8, 6))
    fig.patch.set_facecolor(HEX['black']); ax.set_facecolor(HEX['black'])
    ax.set_aspect('equal')
    
    # Center hub
    hub = plt.Circle((0,0), 0.55, color=HEX['teal'], zorder=5)
    ax.add_patch(hub)
    ax.text(0, 0.12, 'Brezi OS', ha='center', va='center',
            fontsize=12, color='white', fontweight='bold', zorder=6)
    ax.text(0, -0.15, 'BreziIntel', ha='center', va='center',
            fontsize=8, color=HEX['black'], zorder=6)

    products = [
        ('Brezi Wake',    'Cold Brew\nMachine',         0,    2.0,  HEX['teal'],  '06:30'),
        ('Brezi Fuel',    'Smart\nBlender',             1.73, 1.0,  HEX['cyan2'], '07:30'),
        ('Brezi Adapt',   'Adaptogen\nExtractor',       1.73,-1.0,  HEX['teal'],  '07:45'),
        ('Brezi Restore', 'Recovery\nDrink',            0,   -2.0,  HEX['cyan2'], '19:00'),
        ('Brezi Lumen',   'Circadian\nLight',          -1.73,-1.0,  '#555555',    '09:00'),
        ('Brezi Deep',    'Sleep Pad\n(Partner)',       -1.73, 1.0,  '#444444',    '22:00'),
    ]

    for name, desc, x, y, c, time in products:
        # Spoke
        ax.plot([0,x*0.56],[0,y*0.56], color=c, lw=1.2, alpha=0.5, zorder=3)
        # Node
        node = plt.Circle((x, y), 0.45, color=c, zorder=4, alpha=0.9)
        ax.add_patch(node)
        ax.text(x, y+0.12, name, ha='center', va='center',
                fontsize=8, color='white', fontweight='bold', zorder=5)
        ax.text(x, y-0.12, desc, ha='center', va='center',
                fontsize=6.5, color='white', alpha=0.85, zorder=5)
        # Time badge
        ax.text(x, y-0.5, time, ha='center', va='center',
                fontsize=7, color=HEX['teal'], zorder=5)
    
    # Marketplace ring hint
    circle_outer = plt.Circle((0,0), 2.7, fill=False, color='#2BBDBA', 
                               lw=1, linestyle='--', alpha=0.3)
    ax.add_patch(circle_outer)
    ax.text(0, 2.85, 'Brezi Supply Marketplace', ha='center', fontsize=8,
            color='#555555', style='italic')

    ax.set_xlim(-3.2, 3.2); ax.set_ylim(-3.2, 3.2)
    ax.axis('off')
    plt.tight_layout(); return fig

def ch_flywheel():
    """Data flywheel compounding diagram"""
    fig, ax = plt.subplots(figsize=(7, 5))
    fig.patch.set_facecolor(HEX['black']); ax.set_facecolor(HEX['black'])
    ax.set_aspect('equal')

    steps = [
        (0, 1.8,  'More Devices'),
        (1.7, 0.55,'More Signal\nDimensions'),
        (1.7,-1.1, 'Better ML\nAccuracy'),
        (0,  -1.9,'Personalized\nActuation'),
        (-1.7,-1.1,'Better User\nOutcomes'),
        (-1.7, 0.55,'More Users\n& Data'),
    ]
    colors = [HEX['teal'], HEX['teal'], HEX['cyan2'],
              HEX['cyan2'], HEX['teal'], HEX['teal']]

    for i, ((x,y,label), c) in enumerate(zip(steps, colors)):
        node = plt.Circle((x,y), 0.55, color=c, alpha=0.85, zorder=4)
        ax.add_patch(node)
        ax.text(x, y, label, ha='center', va='center',
                fontsize=8.5, color='white', fontweight='bold', zorder=5,
                multialignment='center')
        # Arrow to next
        nx, ny, _ = steps[(i+1) % len(steps)]
        dx, dy = nx-x, ny-y; length=np.sqrt(dx**2+dy**2)
        ux, uy = dx/length, dy/length
        ax.annotate('', xy=(x+ux*0.62, y+uy*0.62),
                    xytext=(x+ux*1.38, y+uy*1.38),
                    arrowprops=dict(arrowstyle='<-', color=HEX['teal'], lw=1.5))

    ax.text(0, 0, '🔄', ha='center', va='center', fontsize=22, zorder=6)
    ax.set_xlim(-2.8, 2.8); ax.set_ylim(-2.8, 2.8)
    ax.axis('off')
    plt.tight_layout(); return fig

def ch_competitive():
    fig, ax = plt.subplots(figsize=(7, 5.5))
    fig.patch.set_facecolor(HEX['grey']); ax.set_facecolor(HEX['grey'])
    ax.axhline(0, color='#CCCCCC', lw=1); ax.axvline(0, color='#CCCCCC', lw=1)
    brands = [
        ('Nespresso',     -0.78,-0.72, HEX['mono']),
        ('Keurig',        -0.88,-0.82, HEX['mono']),
        ('Fellow',        -0.32,-0.42, HEX['mono']),
        ('Plunge',         0.72,-0.22, HEX['mono']),
        ('Oura/WHOOP',     0.68, 0.52, HEX['mono']),
        ('AG1',           -0.12, 0.62, HEX['mono']),
        ('Therabody',      0.48,-0.12, HEX['mono']),
        ('Eight Sleep',    0.58, 0.30, HEX['mono']),
        ('Ice Bath Club',  0.82, 0.18, HEX['mono']),
    ]
    for name, x, y, c in brands:
        ax.scatter(x, y, s=80, color=c, zorder=3)
        ax.annotate(name, (x,y), xytext=(5,4), textcoords='offset points',
                    fontsize=8, color=HEX['black'])
    ax.scatter(0.9, 0.9, s=250, color=HEX['teal'], zorder=6)
    ax.annotate('bre:zi\nOS', (0.9,0.9), xytext=(8,4),
                textcoords='offset points', fontsize=11,
                color=HEX['teal'], fontweight='bold')
    ax.text(0.55, 1.05, '← WHITE SPACE', fontsize=9,
            color=HEX['teal'], fontweight='bold', va='top')
    ax.set_xlim(-1.1,1.1); ax.set_ylim(-1.1,1.1)
    ax.set_xlabel('← Coffee Only     |     Cold Wellness Platform →',
                  color=HEX['mono'], fontsize=9)
    ax.set_ylabel('← Single Product    |    Connected Ecosystem →',
                  color=HEX['mono'], fontsize=9)
    ax.tick_params(bottom=False, left=False, labelbottom=False, labelleft=False)
    for sp in ax.spines.values(): sp.set_visible(False)
    plt.tight_layout(); return fig

def ch_unit_econ():
    fig, ax = plt.subplots(figsize=(6, 3))
    fig.patch.set_facecolor(HEX['black']); ax.set_facecolor(HEX['black'])
    cats   = ['Hardware\nGM', 'Protocol Pack\nGM (Target)', 'Marketplace\nGM (Target)', 'Peloton Sub\nGM (Benchmark)']
    vals   = [40, 63, 72, 69]
    colors = ['#555555', HEX['teal'], HEX['teal'], HEX['cyan2']]
    bars   = ax.bar(cats, vals, color=colors, width=0.5)
    for b, v in zip(bars, vals):
        ax.text(b.get_x()+b.get_width()/2, b.get_height()+1,
                f'{v}%', ha='center', color='white', fontsize=12, fontweight='bold')
    ax.set_ylim(0, 88)
    ax.set_ylabel('Gross Margin %', color='#AAAAAA', fontsize=9)
    ax.tick_params(colors='#AAAAAA', labelsize=9)
    for sp in ax.spines.values(): sp.set_color('#333333')
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_: f'{x:.0f}%'))
    plt.tight_layout(); return fig

def ch_revenue():
    fig, ax = plt.subplots(figsize=(9, 3.5))
    fig.patch.set_facecolor(HEX['black']); ax.set_facecolor(HEX['black'])
    months = list(range(1,19))
    def rev(attach):
        hw  = [399 * min(m*833, 10000)/1000 for m in months]
        sub = []; base = 0
        for m in months:
            new_sub = min(833, 10000-(m-1)*833) * attach
            base = base*0.984 + new_sub
            sub.append(base*25)
        return [h*1000+s for h,s in zip(hw,sub)]
    c1,c2,c3 = rev(0.40), rev(0.50), rev(0.60)
    ax.fill_between(months,[x/1000 for x in c1],[x/1000 for x in c3],
                    color=HEX['teal'], alpha=0.12)
    ax.plot(months,[x/1000 for x in c1],'--',color='#AAAAAA',lw=1.5,label='Conservative 40% attach')
    ax.plot(months,[x/1000 for x in c2],'-', color=HEX['teal'], lw=2.5,label='Base Case 50% attach')
    ax.plot(months,[x/1000 for x in c3],'-', color='white',    lw=1.5,label='Strong Case 60% attach')
    ax.set_xlabel('Month', color='#AAAAAA', fontsize=9)
    ax.set_ylabel('Cumulative Revenue ($K)', color='#AAAAAA', fontsize=9)
    ax.legend(fontsize=8, facecolor='#181818', labelcolor='white', framealpha=0.8)
    ax.tick_params(colors='#AAAAAA', labelsize=8)
    for sp in ax.spines.values(): sp.set_color('#333333')
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_: f'${x:.0f}K'))
    plt.tight_layout(); return fig

def ch_roadmap():
    fig, ax = plt.subplots(figsize=(13, 4.8))
    fig.patch.set_facecolor(HEX['black']); ax.set_facecolor(HEX['black'])
    items = [
        ('Brezi Wake (PrecisionBrew™) — v1 Hardware',    0,  6, HEX['teal'],  5),
        ('iOS App — Manual Protocol Mode',               1,  4, HEX['cyan2'], 4),
        ('Apple HealthKit Integration',                  5,  2, HEX['teal'],  4),
        ('Cold Protocol GTM Campaign — HK/SG',           3,  3, '#DDDDDD',    3),
        ('Brezi Supply Marketplace — MVP',               5,  3, HEX['teal'],  3),
        ('Oura Ring API Integration',                    7,  3, HEX['cyan2'], 4),
        ('University Research Partnership',              2, 10, '#555555',    3),
        ('Brezi Adapt + Brezi Restore — Launch',         9,  4, HEX['teal'],  5),
        ('Brezi Fuel — Launch',                         12,  4, HEX['cyan2'], 5),
        ('WHOOP + Ultrahuman Integration',              12,  4, HEX['teal'],  4),
        ('Brezi Hub — Matter/Thread Controller',        15,  4, HEX['cyan2'], 3),
        ('Brezi Lumen — Software (Philips Hue API)',     9,  4, '#555555',    4),
        ('Series A Fundraise',                          15,  4, HEX['white'], 2),
        ('Brezi Wake v2 — Full Biometric Adaptive',     18,  6, HEX['teal'],  5),
        ('Brezi Lumen — Own Hardware',                  24, 12, '#444444',    4),
        ('BreziOS — Full Ecosystem Platform',           24, 12, HEX['teal'],  3),
    ]
    for label, start, dur, color, row in items:
        ax.barh(row, dur, left=start, color=color, height=0.48, alpha=0.88)
        tc = HEX['black'] if color == HEX['white'] else 'white'
        ax.text(start+0.2, row, label, va='center', fontsize=7,
                color=tc, fontweight='bold' if color==HEX['teal'] else 'normal')
    for x, label, c in [(0,'Phase 1  Build & Launch','#111122'),
                         (9,'Phase 2  Scale & Integrate','#0F1A1A'),
                         (18,'Phase 3  Platform OS','#0A1A10')]:
        ax.axvspan(x, x+9, alpha=0.25, color=c, zorder=0)
        ax.text(x+0.3, 1.5, label, fontsize=8.5, color='#666666', fontweight='bold')
    ax.set_xlim(0,37); ax.set_ylim(1,6.3)
    ax.set_xlabel('Month', color='#AAAAAA', fontsize=9)
    ax.set_yticks([]); ax.tick_params(colors='#AAAAAA', labelsize=8)
    ax.set_xticks([0,3,6,9,12,15,18,21,24,27,30,33,36])
    ax.set_xticklabels([f'M{m}' for m in [0,3,6,9,12,15,18,21,24,27,30,33,36]],
                        fontsize=7.5, color='#AAAAAA')
    for sp in ax.spines.values(): sp.set_color('#333333')
    plt.tight_layout(); return fig

def ch_funds():
    fig, ax = plt.subplots(figsize=(5, 4))
    fig.patch.set_facecolor(HEX['black'])
    labels = ['Hardware\nProduction 35%', 'App & Platform\nDev 25%',
              'GTM & Brand\n20%', 'R&D & Patents\n12%', 'Ops & Team\n8%']
    sizes  = [35, 25, 20, 12, 8]
    colors = [HEX['teal'], HEX['cyan2'], '#4ECFCC', '#666666', '#AAAAAA']
    wedges, texts = ax.pie(sizes, labels=labels, colors=colors, startangle=90,
                           wedgeprops={'width':0.55},
                           textprops={'color':'white', 'fontsize':8.5})
    ax.set_aspect('equal')
    plt.tight_layout(); return fig

def ch_wearables():
    fig, ax = plt.subplots(figsize=(11, 2.8))
    fig.patch.set_facecolor(HEX['black']); ax.set_facecolor(HEX['black'])
    phases = [
        ('PHASE 1 · LAUNCH','Apple\nHealthKit','1B+ users\nZero friction · Day 1', HEX['teal']),
        ('PHASE 2 · M6',    'Oura Ring',       'Best HRV/sleep\nFree OAuth2',       HEX['cyan2']),
        ('PHASE 3 · M12',   'WHOOP',           '300% engagement\nlift benchmark',   '#555555'),
        ('PHASE 4 · M18',   'Ultrahuman\n(APAC)','CGM metabolic\ndata · Partner',  '#444444'),
    ]
    for i,(phase,platform,desc,c) in enumerate(phases):
        x = i*3.0
        rect = mpatches.FancyBboxPatch((x,0.05), 2.65, 2.5,
               boxstyle='round,pad=0.1', fc=c, ec='none', alpha=0.88)
        ax.add_patch(rect)
        ax.text(x+1.32, 2.15, phase,     ha='center', fontsize=7.5,  color='white', fontweight='bold')
        ax.text(x+1.32, 1.45, platform,  ha='center', fontsize=10,   color='white', fontweight='bold')
        ax.text(x+1.32, 0.7,  desc,      ha='center', fontsize=7.5,  color='white', alpha=0.9)
        if i < 3:
            ax.annotate('', xy=(x+2.75,1.25), xytext=(x+2.65,1.25),
                        arrowprops=dict(arrowstyle='->', color=HEX['teal'], lw=1.5))
    ax.set_xlim(-0.3,12.3); ax.set_ylim(0,2.7); ax.axis('off')
    plt.tight_layout(); return fig

# ════════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════════

# S1 COVER
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
tb(s,'12:00 · SAI YING PUN · HONG KONG · 2026',0.6,0.3,14,0.4,size=9,color=MONO)
tb(s,'bre:zi',1.5,1.8,13,3.0,size=120,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
tb(s,'Brewed for your biology.',1.5,5.3,13,0.9,size=30,color=TEAL,align=PP_ALIGN.CENTER)
tb(s,'Venture Capital Pitch Deck  ·  Confidential  ·  2026',1.5,6.4,13,0.5,size=13,color=MONO,align=PP_ALIGN.CENTER)
tb(s,'brezicoffee.com  ·  jackson@brezi.com  ·  Hong Kong',1.5,7.0,13,0.4,size=11,color=MONO,align=PP_ALIGN.CENTER)

# S2 THE VISION
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
secnum(s,'01')
tb(s,'THE VISION',0.6,0.7,14,0.4,size=11,color=TEAL)
tb(s,'We Are Not Building\na Coffee Machine.',0.6,1.2,14.5,2.5,size=54,bold=True,color=WHITE)
tb(s,'We are building the operating system for the human morning —\nstarting with the one ritual that no one has ever optimized.',0.6,4.0,13.5,1.4,size=22,color=MONO)
tb(s,'"Think Apple HomeKit for wellness — but instead of controlling your lights,\nit controls your biology."',0.6,5.7,13.5,1.4,size=19,color=TEAL)
tb(s,'Cold brew is the entry point.  The protocol is the product.  The data is the moat.',0.6,7.3,13.5,0.7,size=16,color=MONO)

# S3 PROBLEM
s = prs.slides.add_slide(BLANK); bg(s, GREY)
eyebrow(s,'THE PROBLEM')
h1(s,'Everything Is Optimized. Except\nthe First Thing You Consume.',size=38,color=BLACK)
items = [
    ('📱','Sleep tracked by Oura. Recovery scored by WHOOP. Zone 2 logged by Garmin.'),
    ('🧊','Cold plunges adopted. Morning light protocols set. Supplements sorted.'),
    ('☕','Yet coffee — the most-repeated morning ritual — is made with tools designed for convenience, not protocol.'),
    ('❌','No product closes the loop: cold exposure + morning consumption + biometric feedback in one system.'),
    ('💡','Cold wellness is a $330M → $660M hardware market. 538M wearables shipped in 2024. The intersection is empty.'),
]
for i,(icon,line) in enumerate(items):
    bar(s, 0.55, 2.2+i*1.1, 0.85)
    tb(s,f'{icon}  {line}',0.75,2.2+i*1.1,14.8,0.95,size=16,color=BLACK,wrap=True)

# S4 MARKET
s = prs.slides.add_slide(BLANK); bg(s, GREY)
eyebrow(s,'MARKET OPPORTUNITY')
h1(s,'Three $B+ Markets Converging.\nBrezi Owns the Intersection.',size=38,color=BLACK)
embed(s, ch_tam(), 0.4, 2.3, 9.5, 3.8)
stats = [('$87B','Wearables market 2025'),('$660M','Cold wellness hardware 2033'),
         ('$40.5B','Coffee consumables 2024'),('32%','Wearable owners with\nattached subscription'),
         ('92%','Would pay more for\nhealth-related features')]
tb(s,'KEY SIGNALS',10.5,2.2,5.2,0.4,size=10,color=TEAL)
for i,(n,l) in enumerate(stats):
    tb(s,n,10.5,2.75+i*1.1,2.1,0.5,size=22,bold=True,color=TEAL)
    tb(s,l,12.7,2.75+i*1.1,2.9,0.55,size=12,color=MONO,wrap=True)

# S5 WHY NOW
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
secnum(s,'01')
eyebrow(s,'WHY NOW',y=0.65)
h1(s,'Three Trends. One Convergence Point.',size=44,color=WHITE)
embed(s, ch_why_now(), 0.5, 2.2, 8.5, 4.2)
tb(s,'CONVERGENCE SIGNALS',9.5,2.2,6.2,0.4,size=10,color=TEAL)
signals = [
    'Cold wellness crossing from sport → daily home practice (81% commercial, 19% residential → flipping)',
    'Wearable subscriptions normalized — 32% attach rate, 92% pay more for health features',
    'Biohacking mainstream: Huberman Lab 1M+ newsletter · Attia 100M+ podcast downloads',
    '"The Ice Bath Club" (SG/HK) already bundles coffee + cold therapy — the pairing is emergent',
    'Matter protocol + Thread mesh = platform-ready IoT infrastructure available today',
]
for i,sig in enumerate(signals):
    bar(s, 9.5, 2.8+i*1.05, 0.85)
    tb(s,sig,9.7,2.8+i*1.05,6.2,0.9,size=12.5,color=WHITE,wrap=True)

# S6 SOLUTION
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
eyebrow(s,'THE SOLUTION')
tb(s,'Introducing Brezi OS',0.6,0.7,14.5,1.1,size=50,bold=True,color=WHITE)
tb(s,'The Connected Wellness Platform for the Protocol Optimizer',0.6,1.9,14,0.7,size=22,color=TEAL)
cols = [('SENSE','Wearables feed biometric state\ninto Brezi OS via OAuth —\nHRV, sleep, recovery, glucose.'),
        ('ANALYZE','BreziIntel determines what\nyour body needs today and\nwhen. Zero guesswork.'),
        ('ACTUATE','The right Brezi device triggers\nat the right time. Cold brew at 9am.\nSmoothie at 1pm. Recovery at 7pm.'),
        ('LEARN','Every behavior is logged.\nDay-over-day comparison.\nThe system compounds over time.')]
for i,(head,body) in enumerate(cols):
    x = 0.4 + i*3.9
    rect = s.shapes.add_shape(1,Inches(x),Inches(3.0),Inches(3.7),Inches(5.5))
    rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(18,18,18)
    rect.line.color.rgb = TEAL
    tb(s,head,x+0.15,3.15,3.4,0.7,size=18,bold=True,color=TEAL)
    tb(s,body,x+0.15,4.0, 3.4,4.0,size=15,color=WHITE,wrap=True)

# S7 EMOTIONAL ARC — THE CENTERPIECE
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
eyebrow(s,'THE PLATFORM VISION')
h1(s,'One Platform. Every Moment That Matters.',size=40,color=WHITE)
embed(s, ch_arc(), 0.3, 2.2, 15.4, 5.5)
tb(s,'"Brezi doesn\'t interrupt your day. It runs underneath it."',0.6,7.85,14.8,0.55,size=13,color=MONO,align=PP_ALIGN.CENTER)

# S8 PRODUCT ECOSYSTEM
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
secnum(s,'02')
eyebrow(s,'THE PRODUCT ECOSYSTEM',y=0.65)
h1(s,'Six Products. One Platform. One Data Loop.',size=40,color=WHITE)
embed(s, ch_ecosystem(), 0.2, 1.8, 8.2, 6.8)
products = [
    ('Brezi Wake',   'NOW',  'Cold brew machine — the daily ritual anchor'),
    ('Brezi Adapt',  '18MO', 'Adaptogen/nootropic cold extractor'),
    ('Brezi Restore','18MO', 'Recovery drink dispenser'),
    ('Brezi Fuel',   '18MO', 'Smart cold nutrition blender'),
    ('Brezi Lumen',  '6MO',  'Circadian light (software-first, Philips Hue API)'),
    ('Brezi Deep',   '36MO', 'Sleep pad (partner Eight Sleep; own hardware later)'),
    ('Brezi Supply', 'NOW',  'Marketplace — Brezi Certified ingredients & supplements'),
]
tb(s,'PRODUCT LINE',8.7,1.8,7,0.4,size=10,color=TEAL)
for i,(name,timing,desc) in enumerate(products):
    tc = TEAL if timing=='NOW' else (WHITE if timing=='18MO' else MONO)
    tb(s,f'{name}',8.7,2.35+i*0.88,2.5,0.45,size=14,bold=True,color=tc)
    tb(s,timing,11.3,2.35+i*0.88,1.2,0.45,size=11,color=MONO)
    tb(s,desc,12.6,2.35+i*0.88,3.0,0.45,size=12,color=WHITE,wrap=True)

# S9 BREZI WAKE — HERO PRODUCT
s = prs.slides.add_slide(BLANK); bg(s, GREY)
eyebrow(s,'HERO PRODUCT')
h1(s,'Brezi Wake — The 12-Minute Protocol',size=40,color=BLACK)
bullets = [
    ('🧊','Thermoelectric cold extraction — no ice, no heat, no compromise. Brews at precisely controlled cold temperature.'),
    ('⏱','12 minutes from cold water to concentrate. Long enough to build habit. Short enough to fit any morning.'),
    ('💧','67% less acidic than hot coffee — measurably better for your gut, cortisol response, and HRV stability.'),
    ('📡','NFC cartridge system: every Brezi Certified Protocol Pack is formula-validated and data-encoded ($0.08/unit).'),
    ('🔵','One cold-blue indicator light makes one slow revolution. No notifications. Just the light. Moving.'),
    ('💰','Hardware $388–$399  |  Protocol Packs $20–35/month  |  Brezi Supply auto-reorder'),
]
for i,(icon,line) in enumerate(bullets):
    bar(s, 0.55, 2.2+i*1.0, 0.75)
    tb(s,f'{icon}  {line}',0.75,2.2+i*1.0,9.8,0.85,size=15.5,color=BLACK,wrap=True)
rect = s.shapes.add_shape(1,Inches(11.0),Inches(1.5),Inches(4.6),Inches(7.0))
rect.fill.solid(); rect.fill.fore_color.rgb = BLACK; rect.line.fill.background()
tb(s,'12:00',10.9,2.8,4.8,2.2,size=82,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
tb(s,'THE BREW WINDOW',10.9,5.2,4.8,0.5,size=10,color=TEAL,align=PP_ALIGN.CENTER)
tb(s,'One revolution.\nOne cold-blue light.\nYour 12 minutes.',10.9,5.9,4.8,1.8,size=13,color=MONO,align=PP_ALIGN.CENTER)

# S10 PLATFORM ARCHITECTURE
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
eyebrow(s,'TECHNOLOGY — BREZI OS')
h1(s,'Every Brew Is a Data Event.\nEvery Event Makes the Next Brew Smarter.',size=36,color=WHITE)
layers = [
    ('LAYER 1  DEVICE','Brezi Wake, Adapt, Restore, Fuel — all on Matter/Thread mesh.\nEach device captures: formula NFC, extraction curve, timestamp.'),
    ('LAYER 2  HEALTH GRAPH','Brezi ingests biometric context from wearables (HRV, sleep, recovery, glucose).\nLinks each device event to your biological state.'),
    ('LAYER 3  BREZINTEL','On-device + cloud ML surfaces adaptive actuation.\n"HRV down 18% — trigger Brezi Wake Recovery Protocol at 07:15."'),
    ('LAYER 4  BREZI SUPPLY','Marketplace auto-replenishment. Brezi Certified partners.\n15–30% take rate. 70%+ GM on consumables.'),
]
for i,(head,body) in enumerate(layers):
    x = 0.5; y = 2.1 + i*1.65
    rect = s.shapes.add_shape(1,Inches(x),Inches(y),Inches(15),Inches(1.4))
    c = RGBColor(43,189,186) if i==0 else (RGBColor(26,138,136) if i==1 else RGBColor(18,18,18))
    rect.fill.solid(); rect.fill.fore_color.rgb = c; rect.line.fill.background()
    tc = BLACK if i<2 else WHITE
    tb(s,head,0.7,y+0.05,5.5,0.5,size=12,bold=True,color=tc)
    tb(s,body,6.3,y+0.05,9.0,1.1,size=13,color=tc,wrap=True)

# S11 DATA FLYWHEEL
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
eyebrow(s,'THE DATA MOAT')
h1(s,'The Dataset That Cannot\nBe Reverse-Engineered.',size=44,color=WHITE)
embed(s, ch_flywheel(), 0.5, 2.0, 7.8, 6.5)
tb(s,'THE TRIPLE LOCK',9.2,2.0,6.5,0.4,size=10,color=TEAL)
moat = [
    ('Unique dataset','Biometric state × brew params × 2h outcomes.\nNo wearable, no coffee brand, no pharma has all three.'),
    ('NFC cartridge auth','Clones copy the machine. They cannot copy the data layer.'),
    ('On-device ML','Your health data never leaves your device without consent.\nPrivacy-first architecture.'),
    ('University research','HKUST/HKU partnership → peer-reviewed publications.\nScientific credibility no competitor can buy.'),
    ('Patent portfolio','Two filings: thermal profiling + biometric-to-thermal mapping.'),
]
for i,(label,desc) in enumerate(moat):
    bar(s, 9.2, 2.6+i*1.3, 1.1)
    tb(s,label,9.4,2.6+i*1.3,6.3,0.45,size=13,bold=True,color=TEAL)
    tb(s,desc,9.4,3.1+i*1.3,6.3,0.75,size=12,color=WHITE,wrap=True)

# S12 COMPETITIVE
s = prs.slides.add_slide(BLANK); bg(s, GREY)
eyebrow(s,'COMPETITIVE LANDSCAPE')
h1(s,'Everyone Is a Tool OR a Brand.\nNobody Is Both.',size=38,color=BLACK)
embed(s, ch_competitive(), 0.4, 2.0, 9.0, 6.6)
tb(s,'THE BREZI ADVANTAGE',10.0,2.0,5.7,0.4,size=10,color=TEAL)
advs = [
    ('Plunge / Ice Barrel','Cold hardware. No daily consumption. No data loop.'),
    ('Oura / WHOOP','Track your biology. Don\'t act on it.'),
    ('Nespresso / Fellow','Premium coffee. Zero wellness framing.'),
    ('AG1 / Momentous','Protocol nutrition. No hardware anchor.'),
    ('Eight Sleep','Sleep temperature. One device, one use case.'),
    ('bre:zi ✓','Closes all loops. Hardware + platform + data + marketplace.'),
]
for i,(brand,desc) in enumerate(advs):
    c = TEAL if brand.startswith('bre:zi') else BLACK
    tb(s,f'• {brand}: {desc}',10.0,2.65+i*1.1,5.7,0.95,size=13,color=c,wrap=True)

# S13 BUSINESS MODEL
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
secnum(s,'03')
eyebrow(s,'BUSINESS MODEL',y=0.65)
h1(s,'Four Revenue Streams.\nAll Compounding.',size=44,color=WHITE)
streams = [
    ('HARDWARE','$388–$399','Device sales.\nTarget GM: 40–45%'),
    ('PROTOCOL PACKS','$20–35/mo','Cartridge subscription.\nNespresso model.\nTarget GM: 63%+'),
    ('BREZI SUPPLY','15–30%\ntake rate','Certified marketplace.\nMomentous, LMNT,\nThesis Nootropics'),
    ('B2B + SPACES','Custom','Hotels, gyms,\ncorporate wellness.\nFrom Month 18'),
]
for i,(stream,price,desc) in enumerate(streams):
    x = 0.4 + i*3.9
    rect = s.shapes.add_shape(1,Inches(x),Inches(3.3),Inches(3.7),Inches(5.1))
    rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(18,18,18)
    rect.line.color.rgb = TEAL if i<3 else RGBColor(60,60,60)
    tb(s,stream,x+0.15,3.45,3.4,0.5,size=11,bold=True,color=TEAL)
    tb(s,price, x+0.15,3.95,3.4,0.85,size=24,bold=True,color=WHITE)
    tb(s,desc,  x+0.15,4.85,3.4,3.2,size=14,color=MONO,wrap=True)

# S14 MARKETPLACE
s = prs.slides.add_slide(BLANK); bg(s, GREY)
eyebrow(s,'BREZI SUPPLY — MARKETPLACE')
h1(s,'Brezi Certified:\nThe Nespresso Model for Premium Wellness.',size=38,color=BLACK)
tb(s,'Partners formulate for the Brezi platform. Cold-extraction optimized. Biometric protocol matched.\n"Brezi Certified" = the MFi for wellness consumables.',0.6,2.5,14.5,1.2,size=17,color=BLACK)
partners = [('Momentous','Day 1','Leading performance supps.\nProtocol Optimizer\'s brand.'),
            ('LMNT','Day 1','Electrolytes. Massive following\nin exact target demographic.'),
            ('Thesis Nootropics','Day 1','Personalized stacks map\nto Brezi OS learning loop.'),
            ('Beam Dream','6MO','Sleep supplement brand.\nEvening protocol slot.'),
            ('Brezi Daily','Own SKU','Morning blend. The\nAG1 moment.'),
            ('Brezi Night','Own SKU','Magnesium + L-theanine.\nNight formula anchor.')]
for i,(name,timing,desc) in enumerate(partners):
    x = 0.5 + (i%3)*5.15; y = 4.0 + (i//3)*2.5
    rect = s.shapes.add_shape(1,Inches(x),Inches(y),Inches(4.85),Inches(2.2))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(43,189,186) if 'Own' in timing else RGBColor(230,230,230)
    rect.line.fill.background()
    tc = WHITE if 'Own' in timing else BLACK
    tb(s,name,  x+0.12,y+0.1,4.5,0.5,size=15,bold=True,color=tc)
    tb(s,timing,x+0.12,y+0.6,4.5,0.4,size=11,color=TEAL if 'Own' not in timing else WHITE)
    tb(s,desc,  x+0.12,y+1.05,4.5,1.0,size=13,color=tc,wrap=True)

# S15 UNIT ECONOMICS
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
eyebrow(s,'UNIT ECONOMICS')
h1(s,'Hardware Opens the Door.\nProtocol Packs Own the Room.',size=40,color=WHITE)
embed(s, ch_unit_econ(), 0.5, 2.0, 7.5, 4.5)
tb(s,'BENCHMARKS',8.5,2.0,7.2,0.4,size=10,color=TEAL)
bmarks = [
    ('Keurig pods','80.5% of US Coffee net sales ($3.2B of $4B) — consumables ARE the business'),
    ('Peloton sub GM','69.1% vs hardware GM ~13% — the margin wedge that changes the P&L'),
    ('Eight Sleep','100% subscription attach year 1 — forced attach, high expectations'),
    ('Target attach','40% conservative · 60% strong case at launch'),
    ('Brezi target GM','Protocol Packs: 63%+ · Marketplace: 72%+'),
    ('18-month MRR','$99K base case (10K units, 50% attach, $25/mo)'),
]
for i,(lbl,val) in enumerate(bmarks):
    bar(s, 8.5, 2.6+i*1.0, 0.85)
    tb(s,lbl, 8.7,2.6+i*1.0,3.2,0.45,size=12,bold=True,color=TEAL)
    tb(s,val, 8.7,3.1+i*1.0,7.0,0.45,size=12,color=WHITE,wrap=True)

# S16 REVENUE PROJECTIONS
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
eyebrow(s,'REVENUE PROJECTIONS')
h1(s,'18-Month Model — 10,000 Units',size=40,color=WHITE)
embed(s, ch_revenue(), 0.4, 2.0, 10.2, 4.5)
tb(s,'BASE CASE (50% attach)',11.2,2.0,4.5,0.4,size=10,color=TEAL)
ms = [('$388–$399','Hardware ASP'),('$25/mo','Protocol Pack price'),
      ('50%','Subscription attach'),('$5.4M','Total 18-mo revenue'),
      ('$99K','MRR at Month 18'),('1.6%/mo','Target monthly churn'),
      ('82%','12-month retention\n(Peloton benchmark)')]
for i,(n,l) in enumerate(ms):
    tb(s,n, 11.2,2.65+i*0.88,2.0,0.45,size=18,bold=True,color=TEAL)
    tb(s,l, 13.3,2.65+i*0.88,2.5,0.5, size=12,color=MONO,wrap=True)

# S17 GTM
s = prs.slides.add_slide(BLANK); bg(s, GREY)
secnum(s,'03')
eyebrow(s,'GO-TO-MARKET')
h1(s,'Seed 500 Right People in Sai Ying Pun.\nLet Them Pull the Brand Forward.',size=36,color=BLACK)
phases = [
    ('PHASE 1\nDays 1–30\n~$40K HKD',
     '• Announce university research partnership\n• 15–20 Protocol Practitioners seeded\n  (coaches, sports med, Hyrox, CrossFit)\n• "The Cold Protocol" newsletter — Day 1\n• Waitlist page + 90-sec hero film'),
    ('PHASE 2\nDays 31–60\n~$120K HKD',
     '• Invite-only pop-up, Sai Ying Pun HK\n  (200 spots — stone floors, no florals)\n• Live Brezi Wake demo + cold plunge station\n• Hero film "Start Cold" debut\n• Huberman-adjacent podcast ad buy'),
    ('PHASE 3\nDays 61–90\n~$80K HKD',
     '• 30-day Baseline Protocol (not a "challenge")\n• First 500 units at $388 USD\n• Protocol Packs as post-purchase upsell\n• Referral network → 500 becomes 2,000'),
]
for i,(phase,bullets) in enumerate(phases):
    x = 0.4 + i*5.2
    rect = s.shapes.add_shape(1,Inches(x),Inches(2.2),Inches(5.0),Inches(6.3))
    c = RGBColor(43,189,186) if i==0 else (RGBColor(26,138,136) if i==1 else BLACK)
    rect.fill.solid(); rect.fill.fore_color.rgb = c; rect.line.fill.background()
    tc = BLACK if i==0 else WHITE
    tb(s,phase,x+0.15,2.35,4.7,1.1,size=13,bold=True,color=tc)
    tb(s,bullets,x+0.15,3.6,4.7,4.5,size=13,color=tc,wrap=True)
tb(s,'Total: ~$240K HKD (~$30K USD)  ·  KPIs: 2,000 waitlist · 500 paying · 40% Protocol Pack attach',
   0.5,8.5,15,0.45,size=11,color=MONO)

# S18 PRODUCT ROADMAP
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
secnum(s,'03')
eyebrow(s,'PRODUCT ROADMAP',y=0.65)
h1(s,'0 → 36 Months: Device → Platform → OS',size=40,color=WHITE)
embed(s, ch_roadmap(), 0.2, 2.5, 15.6, 6.0)

# S19 WEARABLES
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
eyebrow(s,'WEARABLE INTEGRATIONS')
h1(s,'Every Device You Already Trust.\nNow Connected to Your Morning.',size=40,color=WHITE)
embed(s, ch_wearables(), 0.5, 2.5, 15.0, 4.0)
tb(s,'Cronometer × WHOOP: 300% engagement lift  ·  Strava × Apple Watch: 20% YoY upload increase',
   0.5,7.0,15,0.5,size=12,color=MONO,align=PP_ALIGN.CENTER)

# S20 WHY US
s = prs.slides.add_slide(BLANK); bg(s, GREY)
secnum(s,'04')
eyebrow(s,'WHY US')
h1(s,'We Are the Customer.\nThat Is the Unfair Advantage.',size=40,color=BLACK)
left = [
    ('🎓','Jackson (Founder) — UC Berkeley. First Protocol Optimizer on the cap table.'),
    ('🏭','Engineering team from semiconductor + consumer electronics (DJI/Sony tier).\nThermal precision hardware is our DNA.'),
    ('🏃','Every team member is an active sports and wellness practitioner.\nWe ship what we test on ourselves.'),
    ('🔬','University research partnership — HKUST/HKU Medicine — in active discussions.'),
]
right = [
    ('📍','Hong Kong base: Asia\'s densest high-performance urban cohort.\nOur test market is our backyard.'),
    ('🤝','Graves Ho (Co-Founder): HK market depth, manufacturing network, operations.'),
    ('💡','Semiconductor background = miniaturized thermoelectric design advantage\nthat consumer brands cannot easily replicate.'),
    ('📊','Phase 1 research complete: 7 interviews + 181 backer surveys.\nWe know our customer. We have their data.'),
]
for i,(icon,line) in enumerate(left):
    bar(s, 0.5, 2.3+i*1.5, 1.2)
    tb(s,f'{icon}  {line}',0.7,2.3+i*1.5,7.2,1.2,size=15,color=BLACK,wrap=True)
for i,(icon,line) in enumerate(right):
    bar(s, 8.3, 2.3+i*1.5, 1.2)
    tb(s,f'{icon}  {line}',8.5,2.3+i*1.5,7.2,1.2,size=15,color=BLACK,wrap=True)

# S21 TRACTION
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
eyebrow(s,'TRACTION & VALIDATION')
h1(s,'Not a Concept.\nA Validated Product with Early Signal.',size=40,color=WHITE)
traction = [
    ('✅','Kickstarter campaign successfully funded — proof of market demand for cold brew hardware'),
    ('✅','Live Shopify store (brezicoffee.com) + Taobao presence — two-market distribution already live'),
    ('✅','Phase 1 customer research complete: 7 in-depth interviews + 181 backer surveys'),
    ('✅','iOS technical spec complete · UX copy brief complete · App in active development'),
    ('✅','Roaster partnerships live: Regent Coffee + Stitch Coffee on platform'),
    ('✅','Hardware: Thermoelectric cooling IP validated. PrecisionBrew™ technology proven.'),
    ('✅','Full team assembled: founder + semiconductor engineers + operations'),
    ('🔄','University research partnership: HKUST / HKU Medicine — in active discussions'),
]
for i,(icon,line) in enumerate(traction):
    bar(s, 0.5, 1.7+i*0.9, 0.7)
    tb(s,f'{icon}  {line}',0.7,1.7+i*0.9,14.8,0.8,size=16,color=WHITE,wrap=True)

# S22 THE ASK
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
secnum(s,'04')
eyebrow(s,'THE ASK',y=0.65)
h1(s,'Join the Protocol.',size=52,color=WHITE)
# Capital box
rect = s.shapes.add_shape(1,Inches(0.5),Inches(2.7),Inches(7.3),Inches(5.8))
rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(18,18,18); rect.line.color.rgb = TEAL
tb(s,'CAPITAL',0.7,2.85,6.9,0.4,size=11,bold=True,color=TEAL)
tb(s,'[RAISE AMOUNT]',0.7,3.35,6.9,0.9,size=30,bold=True,color=WHITE)
tb(s,'at [VALUATION] pre-money',0.7,4.35,6.9,0.5,size=16,color=MONO)
for i,item in enumerate([
    'Complete precision batch — 10,000 units',
    'iOS app + platform build (3-month sprint)',
    'Cold Protocol GTM Campaign — HK + SG',
    'University research partnership + 2 patent filings',
    'BreziRest/Adapt R&D initiation',
    'Brezi Supply marketplace MVP']):
    tb(s,f'· {item}',0.7,5.0+i*0.43,6.9,0.4,size=13,color=WHITE)
# Talent box
rect = s.shapes.add_shape(1,Inches(8.3),Inches(2.7),Inches(7.3),Inches(5.8))
rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(18,18,18); rect.line.color.rgb = CYAN2
tb(s,'TALENT — OPEN ROLES',8.5,2.85,6.9,0.4,size=11,bold=True,color=TEAL)
tb(s,'Building the Core Team Now',8.5,3.35,6.9,0.7,size=22,bold=True,color=WHITE)
for i,role in enumerate([
    'Senior thermal systems hardware engineer',
    'iOS engineer — React Native + HealthKit APIs',
    'Product designer — industrial + UX',
    'Head of Protocol Community — HK/SG',
    'Research coordinator — university partnership',
    'Brezi Supply — marketplace partnerships lead']):
    tb(s,f'· {role}',8.5,4.15+i*0.55,6.9,0.5,size=14,color=WHITE)

# S23 USE OF FUNDS
s = prs.slides.add_slide(BLANK); bg(s, GREY)
eyebrow(s,'USE OF FUNDS')
h1(s,'Every Dollar Has a Job.',size=44,color=BLACK)
embed(s, ch_funds(), 0.3, 1.8, 7.5, 6.8)
tb(s,'ALLOCATION RATIONALE',8.5,2.0,7.2,0.4,size=10,color=TEAL)
allocs = [
    ('35%  Hardware Production','First 10,000-unit precision batch. NFC cartridge tooling. Quality assurance.'),
    ('25%  App & Platform Dev','3-month iOS sprint + backend data infrastructure + Apple HealthKit integration.'),
    ('20%  GTM & Brand','90-day Cold Protocol Campaign. SYP pop-up. Hero film. Podcast media buy.'),
    ('12%  R&D & Patents','Two patent filings. BreziRest/Adapt R&D. University research seed funding.'),
    (' 8%  Operations & Team','Critical hires. Legal. Corporate structure. HK incorporation.'),
]
for i,(alloc,detail) in enumerate(allocs):
    bar(s, 8.5, 2.65+i*1.2, 1.0)
    tb(s,alloc, 8.7,2.65+i*1.2,7.0,0.45,size=13,bold=True,color=BLACK)
    tb(s,detail,8.7,3.15+i*1.2,7.0,0.6, size=13,color=MONO,wrap=True)

# S24 CLOSING
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
tb(s,'bre:zi',1.5,1.8,13,2.5,size=100,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
tb(s,'Brewed for your biology.',1.5,4.7,13,0.9,size=28,color=TEAL,align=PP_ALIGN.CENTER)
tb(s,'"The compound interest of your wellness routine."',1.5,5.8,13,0.7,size=18,color=MONO,align=PP_ALIGN.CENTER)
tb(s,'brezicoffee.com  ·  jackson@brezi.com  ·  Hong Kong',1.5,6.9,13,0.5,size=14,color=MONO,align=PP_ALIGN.CENTER)

prs.save(OUT)
print(f'✅  {OUT}  ({len(prs.slides)} slides)')

"""
Brezi VC Pitch Deck Builder
Generates Brezi_VC_Pitch_Deck.pptx with embedded matplotlib charts
"""

import os, io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.ticker as mticker
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ──────────────────────────────────────────────────────────────────
BLACK   = RGBColor(0,   0,   0)
WHITE   = RGBColor(255, 255, 255)
GREY    = RGBColor(245, 245, 245)
BONE    = RGBColor(240, 235, 225)
TEAL    = RGBColor(43,  189, 186)
CYAN2   = RGBColor(26,  138, 136)
MONO    = RGBColor(140, 140, 140)

HEX_BLACK = '#000000'
HEX_WHITE = '#FFFFFF'
HEX_GREY  = '#F5F5F5'
HEX_BONE  = '#F0EBE1'
HEX_TEAL  = '#2BBDBA'
HEX_CYAN2 = '#1A8A88'
HEX_MONO  = '#8C8C8C'

OUT = 'projects/brezi-coffee/Brezi_VC_Pitch_Deck.pptx'
os.makedirs('projects/brezi-coffee', exist_ok=True)

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)

BLANK = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────────

def bg(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb

def tb(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb

def eyebrow(slide, text, y=0.25):
    tb(slide, text, 0.6, y, 14, 0.4, size=10, bold=False, color=TEAL)

def title_line(slide, text, y=0.65, size=44, color=BLACK):
    tb(slide, text, 0.6, y, 14.8, 1.2, size=size, bold=True, color=color)

def accent_bar(slide, x, y, h):
    """2pt Glacier Teal left accent bar"""
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.05), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

def section_number(slide, num):
    tb(slide, num, 0.6, 0.35, 2, 0.5, size=11, bold=False, color=TEAL)

def img_slide(slide, fig, x, y, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='none', transparent=True)
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))
    plt.close(fig)

def bullet_block(slide, items, x, y, w, h, size=17, color=BLACK, gap=0.42):
    for i, (icon, line) in enumerate(items):
        tb(slide, f"{icon}  {line}", x, y + i*gap, w, gap+0.05,
           size=size, color=color, wrap=True)

# ════════════════════════════════════════════════════════════════════════════
# CHART HELPERS
# ════════════════════════════════════════════════════════════════════════════

def chart_tam():
    """TAM / SAM / SOM horizontal bar"""
    fig, ax = plt.subplots(figsize=(8, 3.2))
    fig.patch.set_facecolor(HEX_GREY)
    ax.set_facecolor(HEX_GREY)
    labels = ['TAM\n(Cold Wellness + Wearables\n+ Premium Coffee)', 'SAM\n(Addressable with\nBrezi Platform)', 'SOM\n(3–5 Year\nCapture)']
    vals   = [89.3, 0.694, 0.0135]
    colors = [HEX_CYAN2, HEX_TEAL, '#48DDDA']
    bars = ax.barh(labels, vals, color=colors, height=0.55)
    for bar, val in zip(bars, vals):
        label = f'${val:.1f}B' if val >= 1 else f'${val*1000:.0f}M'
        ax.text(bar.get_width() + 0.5, bar.get_y() + bar.get_height()/2,
                label, va='center', color=HEX_BLACK, fontsize=11, fontweight='bold')
    ax.set_xlim(0, 105)
    ax.set_xlabel('USD Billion', color=HEX_MONO, fontsize=9)
    ax.tick_params(colors=HEX_BLACK, labelsize=9)
    ax.spines[:].set_visible(False)
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'${x:.0f}B'))
    plt.tight_layout()
    return fig

def chart_why_now():
    """3 converging trend lines"""
    fig, ax = plt.subplots(figsize=(7, 3.2))
    fig.patch.set_facecolor(HEX_BLACK)
    ax.set_facecolor(HEX_BLACK)
    years = [2020, 2021, 2022, 2023, 2024, 2025, 2026]
    cold  = [100, 118, 145, 178, 210, 255, 310]
    wear  = [100, 110, 128, 148, 165, 190, 218]
    bioh  = [100, 112, 130, 155, 185, 220, 262]
    ax.plot(years, cold, color=HEX_TEAL,  lw=2.5, marker='o', ms=5, label='Cold Wellness')
    ax.plot(years, wear, color='#AAAAAA', lw=2.5, marker='s', ms=5, label='Wearables')
    ax.plot(years, bioh, color='#FFFFFF', lw=2.5, marker='^', ms=5, label='Biohacking')
    ax.axvline(x=2025, color=HEX_TEAL, lw=1, linestyle='--', alpha=0.6)
    ax.text(2025.1, 270, 'Brezi\nenters', color=HEX_TEAL, fontsize=8)
    ax.set_ylabel('Indexed Growth (2020=100)', color='#AAAAAA', fontsize=8)
    ax.tick_params(colors='#AAAAAA', labelsize=8)
    ax.legend(fontsize=8, facecolor='#111111', labelcolor='white', framealpha=0.8)
    for spine in ax.spines.values(): spine.set_color('#333333')
    plt.tight_layout()
    return fig

def chart_competitive():
    """2x2 competitive quadrant"""
    fig, ax = plt.subplots(figsize=(7, 5.5))
    fig.patch.set_facecolor(HEX_GREY)
    ax.set_facecolor(HEX_GREY)
    ax.axhline(0, color='#CCCCCC', lw=1)
    ax.axvline(0, color='#CCCCCC', lw=1)
    brands = [
        ('Nespresso',    -0.75, -0.70, HEX_MONO),
        ('Keurig',       -0.85, -0.80, HEX_MONO),
        ('Fellow',       -0.30, -0.40, HEX_MONO),
        ('Plunge',        0.70, -0.20, HEX_MONO),
        ('Oura / WHOOP',  0.65,  0.50, HEX_MONO),
        ('AG1',          -0.10,  0.60, HEX_MONO),
        ('Therabody',     0.45, -0.10, HEX_MONO),
        ('Ice Bath Club', 0.80,  0.20, HEX_MONO),
    ]
    for name, x, y, c in brands:
        ax.scatter(x, y, s=90, color=c, zorder=3)
        ax.annotate(name, (x, y), textcoords='offset points', xytext=(6, 4),
                    fontsize=8, color=HEX_BLACK)
    ax.scatter(0.85, 0.85, s=220, color=HEX_TEAL, zorder=5)
    ax.annotate('bre:zi', (0.85, 0.85), textcoords='offset points',
                xytext=(8, 4), fontsize=11, color=HEX_TEAL, fontweight='bold')
    ax.set_xlim(-1.1, 1.1); ax.set_ylim(-1.1, 1.1)
    ax.set_xlabel('← Coffee Only     |     Cold Wellness →',
                  color=HEX_MONO, fontsize=9)
    ax.set_ylabel('← Low Protocol    |    High Protocol →',
                  color=HEX_MONO, fontsize=9)
    ax.tick_params(bottom=False, left=False, labelbottom=False, labelleft=False)
    for spine in ax.spines.values(): spine.set_visible(False)
    ax.text(-1.05, 1.05, 'Biometric\nplatform',   fontsize=8, color='#AAAAAA', va='top')
    ax.text( 0.55, 1.05, 'WHITE SPACE →',         fontsize=9, color=HEX_TEAL,  va='top', fontweight='bold')
    ax.text(-1.05,-1.05, 'Commodity\ncoffee',      fontsize=8, color='#AAAAAA', va='bottom')
    ax.text( 0.55,-1.05, 'Cold hardware\nno ritual', fontsize=8, color='#AAAAAA', va='bottom')
    plt.tight_layout()
    return fig

def chart_unit_economics():
    """Hardware vs subscription margin comparison"""
    fig, ax = plt.subplots(figsize=(6, 3.2))
    fig.patch.set_facecolor(HEX_BLACK)
    ax.set_facecolor(HEX_BLACK)
    cats   = ['Hardware\nGM', 'Subscription\nGM (Target)', 'Peloton\nSub GM\n(Benchmark)']
    vals   = [40, 62, 69]
    colors = ['#555555', HEX_TEAL, HEX_CYAN2]
    bars = ax.bar(cats, vals, color=colors, width=0.5)
    for bar, val in zip(bars, vals):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1,
                f'{val}%', ha='center', color='white', fontsize=12, fontweight='bold')
    ax.set_ylim(0, 85)
    ax.set_ylabel('Gross Margin %', color='#AAAAAA', fontsize=9)
    ax.tick_params(colors='#AAAAAA', labelsize=9)
    for spine in ax.spines.values(): spine.set_color('#333333')
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'{x:.0f}%'))
    plt.tight_layout()
    return fig

def chart_revenue():
    """18-month revenue projection, 3 scenarios"""
    fig, ax = plt.subplots(figsize=(8, 3.5))
    fig.patch.set_facecolor(HEX_BLACK)
    ax.set_facecolor(HEX_BLACK)
    months = list(range(1, 19))
    def rev(attach):
        hw  = [399 * min(m * 833, 10000) / 1000 for m in months]
        sub = []
        base = 0
        for m in months:
            new_sub = min(833, 10000 - (m-1)*833) * attach
            base = base * 0.984 + new_sub
            sub.append(base * 25)
        return [h*1000 + s for h, s in zip(hw, sub)]
    cons = rev(0.40)
    base = rev(0.50)
    aggr = rev(0.60)
    ax.fill_between(months, [c/1000 for c in cons], [a/1000 for a in aggr],
                    color=HEX_TEAL, alpha=0.15)
    ax.plot(months, [c/1000 for c in cons], '--', color='#AAAAAA', lw=1.5, label='Conservative (40% attach)')
    ax.plot(months, [b/1000 for b in base],  '-',  color=HEX_TEAL,  lw=2.5, label='Base Case (50% attach)')
    ax.plot(months, [a/1000 for a in aggr],  '-',  color=HEX_WHITE, lw=1.5, label='Strong Case (60% attach)')
    ax.set_xlabel('Month', color='#AAAAAA', fontsize=9)
    ax.set_ylabel('Cumulative Revenue (USD $K)', color='#AAAAAA', fontsize=9)
    ax.legend(fontsize=8, facecolor='#111111', labelcolor='white', framealpha=0.8)
    ax.tick_params(colors='#AAAAAA', labelsize=8)
    for spine in ax.spines.values(): spine.set_color('#333333')
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'${x:.0f}K'))
    plt.tight_layout()
    return fig

def chart_roadmap():
    """Product roadmap Gantt-style"""
    fig, ax = plt.subplots(figsize=(12, 4.5))
    fig.patch.set_facecolor(HEX_BLACK)
    ax.set_facecolor(HEX_BLACK)
    items = [
        # (label, start_mo, dur_mo, color, row)
        ('PrecisionBrew™ v1 — Hardware',          0,  6, HEX_TEAL,  5),
        ('iOS App v1 (Manual Protocol Mode)',      1,  4, HEX_CYAN2, 4),
        ('Apple HealthKit Integration',            5,  2, HEX_TEAL,  4),
        ('Cold Protocol Campaign — HK Launch',     3,  3, '#DDDDDD', 3),
        ('Oura Ring API Integration',              7,  3, HEX_TEAL,  4),
        ('University Research Partnership',        2, 10, '#666666', 3),
        ('Cartridge Subscription Launch',          5,  2, HEX_CYAN2, 3),
        ('WHOOP Integration',                     12,  4, HEX_TEAL,  4),
        ('BreziRest Pod — R&D',                    9,  9, '#444444', 5),
        ('BreziBar — Concept + B2B Pilots',       24, 12, '#444444', 5),
        ('Series A Fundraise',                    15,  4, HEX_WHITE, 2),
    ]
    yticks, ylabels = [], []
    for label, start, dur, color, row in items:
        ax.barh(row, dur, left=start, color=color, height=0.5, alpha=0.9)
        ax.text(start + 0.2, row, label, va='center', fontsize=7.5,
                color='black' if color in (HEX_WHITE, '#DDDDDD') else 'white',
                fontweight='bold' if color == HEX_TEAL else 'normal')
    # Phase bands
    for x, label, c in [(0,'Phase 1\nBuild & Launch','#1A1A2E'),
                         (9,'Phase 2\nScale & Integrate','#0F1A1A'),
                         (18,'Phase 3\nPlatform','#0A1A10')]:
        ax.axvspan(x, x+9, alpha=0.3, color=c, zorder=0)
        ax.text(x+0.3, 1.3, label, fontsize=8, color='#AAAAAA')
    ax.set_xlim(0, 37)
    ax.set_ylim(0.8, 6.2)
    ax.set_xlabel('Month', color='#AAAAAA', fontsize=9)
    ax.tick_params(colors='#AAAAAA', labelsize=8)
    ax.set_yticks([])
    ax.set_xticks([0,3,6,9,12,15,18,21,24,27,30,33,36])
    ax.set_xticklabels([f'M{m}' for m in [0,3,6,9,12,15,18,21,24,27,30,33,36]],
                        fontsize=7.5, color='#AAAAAA')
    for spine in ax.spines.values(): spine.set_color('#333333')
    plt.tight_layout()
    return fig

def chart_use_of_funds():
    """Donut chart — use of funds"""
    fig, ax = plt.subplots(figsize=(5, 4))
    fig.patch.set_facecolor(HEX_BLACK)
    labels = ['Hardware\nProduction\n35%', 'App & Platform\nDev 25%',
              'GTM & Brand\nActivation 20%', 'R&D &\nPatents 12%',
              'Operations\n& Team 8%']
    sizes  = [35, 25, 20, 12, 8]
    colors = [HEX_TEAL, HEX_CYAN2, '#4ECFCC', '#777777', '#AAAAAA']
    wedges, texts = ax.pie(sizes, labels=labels, colors=colors,
                           startangle=90, wedgeprops={'width': 0.55},
                           textprops={'color': 'white', 'fontsize': 8})
    ax.set_aspect('equal')
    plt.tight_layout()
    return fig

def chart_wearable_phases():
    """Wearable integration phase visual"""
    fig, ax = plt.subplots(figsize=(10, 2.8))
    fig.patch.set_facecolor(HEX_BLACK)
    ax.set_facecolor(HEX_BLACK)
    phases = [
        ('PHASE 1\nLaunch', 'Apple\nHealthKit', '1B+ users\nZero friction\nDay 1', 0),
        ('PHASE 2\nMonth 6', 'Oura Ring', 'Best HRV data\nFree OAuth2\nSleep gold standard', 3),
        ('PHASE 3\nMonth 12', 'WHOOP', 'Recovery protocol\nusers\n300% engagement lift', 6),
        ('PHASE 4\nMonth 18', 'Ultrahuman\n(APAC)', 'CGM metabolic\ndata\nInvite-only partner', 9),
    ]
    for i, (phase, platform, desc, x) in enumerate(phases):
        c = HEX_TEAL if i == 0 else (HEX_CYAN2 if i == 1 else '#555555')
        rect = mpatches.FancyBboxPatch((x, 0.1), 2.5, 2.4,
               boxstyle='round,pad=0.1', fc=c, ec='none', alpha=0.85)
        ax.add_patch(rect)
        ax.text(x+1.25, 2.0, phase, ha='center', va='center',
                fontsize=8, color='white', fontweight='bold')
        ax.text(x+1.25, 1.35, platform, ha='center', va='center',
                fontsize=9.5, color='white', fontweight='bold')
        ax.text(x+1.25, 0.55, desc, ha='center', va='center',
                fontsize=7.5, color='white', alpha=0.9)
        if i < 3:
            ax.annotate('', xy=(x+2.6, 1.25), xytext=(x+2.5, 1.25),
                        arrowprops=dict(arrowstyle='->', color=HEX_TEAL, lw=1.5))
    ax.set_xlim(-0.3, 12); ax.set_ylim(0, 2.6)
    ax.axis('off')
    plt.tight_layout()
    return fig

# ════════════════════════════════════════════════════════════════════════════
# BUILD SLIDES
# ════════════════════════════════════════════════════════════════════════════

# ── S1: COVER ────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
tb(s, '12:00', 0.6, 0.3, 4, 0.5, size=10, color=MONO)
tb(s, 'bre:zi', 2, 2.2, 12, 3.5, size=120, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 'Brewed for your biology.', 2, 5.5, 12, 1, size=32, bold=False, color=TEAL, align=PP_ALIGN.CENTER)
tb(s, 'Venture Capital Pitch Deck · Confidential · 2026', 2, 6.6, 12, 0.5, size=13, color=MONO, align=PP_ALIGN.CENTER)
tb(s, 'brezicoffee.com · jackson@brezi.com', 2, 7.3, 12, 0.5, size=11, color=MONO, align=PP_ALIGN.CENTER)

# ── S2: THE VISION ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
section_number(s, '01')
tb(s, 'THE VISION', 0.6, 0.75, 14, 0.5, size=11, color=TEAL)
tb(s, 'We Are Not Building\na Coffee Machine.', 0.6, 1.3, 14, 2.2, size=56, bold=True, color=WHITE)
tb(s, 'We are building the operating system for the human morning —\nstarting with the one ritual no one has optimized yet.', 0.6, 4.0, 14, 1.5, size=22, color=MONO)
tb(s, 'Cold brew is the entry point.\nThe protocol is the product.\nThe data is the moat.', 0.6, 5.7, 14, 2.0, size=20, color=TEAL)

# ── S3: PROBLEM ──────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, GREY)
eyebrow(s, 'THE PROBLEM')
title_line(s, 'Everything Is Optimized.\nExcept the First Thing They Consume.', size=38)
items = [
    ('📱', 'Sleep tracked by Oura. Recovery scored by WHOOP. Training logged by Garmin.'),
    ('🧊', 'Cold plunges adopted. Morning light protocols running. Zone 2 scheduled.'),
    ('☕', 'Yet coffee — the most repeated morning ritual — is made with tools designed for convenience, not protocol.'),
    ('❌', 'No product connects cold exposure + morning consumption + biometric feedback in a single loop.'),
    ('📊', 'Cold plunge market: $330M → $660M by 2033 (8.1% CAGR). 538M wearables shipped in 2024. Nobody owns the intersection.'),
]
bullet_block(s, items, 0.6, 2.2, 14.8, 5.5, size=16, color=BLACK, gap=0.95)

# ── S4: MARKET OPPORTUNITY ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, GREY)
eyebrow(s, 'MARKET OPPORTUNITY')
title_line(s, 'Three $B+ Markets Converging.\nBrezi Owns the Intersection.', size=38, color=BLACK)
fig = chart_tam()
img_slide(s, fig, 0.5, 2.3, 9.5, 4.0)
# Side stats
tb(s, 'KEY SIGNALS', 10.5, 2.3, 5, 0.5, size=10, color=TEAL)
stats = [
    ('$87B', 'Wearables market 2025'),
    ('$660M', 'Cold wellness hardware 2033'),
    ('$40.5B', 'Coffee pod/consumables 2024'),
    ('32%', 'Wearable owners have\nan attached subscription'),
    ('92%', 'Would pay more for\nhealth features'),
]
for i, (num, label) in enumerate(stats):
    tb(s, num,   10.5, 2.9 + i*1.1, 2.2, 0.5, size=22, bold=True, color=TEAL)
    tb(s, label, 12.8, 2.9 + i*1.1, 2.7, 0.5, size=12, color=MONO, wrap=True)

# ── S5: WHY NOW ───────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
section_number(s, '02')
eyebrow(s, 'WHY NOW', y=0.6)
title_line(s, 'Three Trends. One Convergence Point.', size=44, color=WHITE)
fig = chart_why_now()
img_slide(s, fig, 0.5, 2.2, 8.5, 4.2)
tb(s, 'INFLECTION SIGNALS', 9.5, 2.2, 6, 0.5, size=10, color=TEAL)
signals = [
    ('Cold wellness crossing from sport\ninto daily home practice'),
    ('Wearable subscriptions normalized —\n32% attach rate (Parks Assoc.)'),
    ('Biohacking cultural mainstream:\nHuberman Lab 1M+ newsletter subs'),
    ('"The Ice Bath Club" (SG/HK)\nalready bundles coffee + cold therapy'),
    ('API access to Oura, WHOOP, Apple\nHealth now mature and open'),
]
for i, sig in enumerate(signals):
    accent_bar(s, 9.5, 2.85 + i*1.0, 0.7)
    tb(s, sig, 9.7, 2.85 + i*1.0, 6, 0.9, size=13, color=WHITE, wrap=True)

# ── S6: SOLUTION ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
eyebrow(s, 'THE SOLUTION', y=0.3)
tb(s, 'Introducing bre:zi', 0.6, 0.75, 14, 1.2, size=52, bold=True, color=WHITE)
tb(s, 'The Cold Protocol Platform', 0.6, 1.9, 14, 0.7, size=26, color=TEAL)
# 3-column how it works
cols = [
    ('READS', 'Your biometric state.\nHRV, sleep score,\nrecovery index —\nbefore every brew.'),
    ('BREWS', 'PrecisionBrew™.\n12 minutes. Cold.\nFormula-matched\nto your biology today.'),
    ('LEARNS', 'Every session logged.\nEvery outcome tracked.\nThe longer you use it,\nthe smarter it gets.'),
]
for i, (head, body) in enumerate(cols):
    x = 0.6 + i * 5.1
    rect = s.shapes.add_shape(1, Inches(x), Inches(3.0), Inches(4.7), Inches(5.0))
    rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(20, 20, 20)
    rect.line.color.rgb = TEAL
    tb(s, head, x+0.2, 3.2, 4.3, 0.7, size=20, bold=True, color=TEAL)
    tb(s, body, x+0.2, 4.0, 4.3, 3.5, size=17, color=WHITE, wrap=True)

# ── S7: PRODUCT ───────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, GREY)
eyebrow(s, 'THE PRODUCT')
title_line(s, 'PrecisionBrew™ — The 12-Minute Protocol', size=38, color=BLACK)
items = [
    ('🧊', 'Thermoelectric cold extraction — no ice, no heat. Brews at controlled cold temperature.'),
    ('⏱', '12 minutes from cold water to perfect cold brew concentrate. The ritual window.'),
    ('💧', '67% less acidic than hot coffee — medically relevant, not just a feature.'),
    ('📡', 'NFC cartridge system: every Protocol Pack is formula-validated and data-encoded ($0.08/unit COGS).'),
    ('📱', 'Companion app: morning state check-in → brew profile recommendation → 2h outcome log.'),
    ('💰', 'Hardware $388–$399 | Protocol Packs $20–35/month subscription'),
]
bullet_block(s, items, 0.6, 2.1, 9.8, 6.0, size=16, color=BLACK, gap=0.95)
# 12:00 graphic panel right
rect = s.shapes.add_shape(1, Inches(10.8), Inches(1.5), Inches(4.8), Inches(7.0))
rect.fill.solid(); rect.fill.fore_color.rgb = BLACK
rect.line.fill.background()
tb(s, '12:00', 10.8, 3.0, 4.8, 2.0, size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 'THE BREW WINDOW', 10.8, 5.2, 4.8, 0.5, size=10, color=TEAL, align=PP_ALIGN.CENTER)
tb(s, 'One slow revolution.\nOne cold-blue indicator light.\nYour 12 minutes.', 10.8, 5.9, 4.8, 1.5, size=13, color=MONO, align=PP_ALIGN.CENTER)

# ── S8: COMPETITIVE LANDSCAPE ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, GREY)
eyebrow(s, 'COMPETITIVE LANDSCAPE')
title_line(s, 'Everyone Is a Tool OR a Brand.\nNobody Is Both.', size=38, color=BLACK)
fig = chart_competitive()
img_slide(s, fig, 0.4, 2.1, 8.8, 6.5)
tb(s, 'WHY BREZI WINS', 10.0, 2.1, 5.7, 0.5, size=11, color=TEAL)
advantages = [
    ('Cold plunges', 'own the cold exposure — no daily consumption ritual, no data loop'),
    ('Nespresso', 'owns the pod system — zero protocol/wellness framing'),
    ('Oura / WHOOP', 'own the biometric data — no ingestible system'),
    ('AG1', 'owns functional nutrition — consumable only, no hardware anchor'),
    ('bre:zi', '✓ Closes all three loops in one morning ritual object'),
]
for i, (brand, desc) in enumerate(advantages):
    c = TEAL if brand == 'bre:zi' else BLACK
    tb(s, f'• {brand}: {desc}', 10.0, 2.75 + i*1.15, 5.7, 1.0,
       size=13, color=c, wrap=True)

# ── S9: BUSINESS MODEL ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
section_number(s, '03')
eyebrow(s, 'BUSINESS MODEL', y=0.6)
title_line(s, 'Hardware Opens the Door.\nProtocol Packs Own the Room.', size=44, color=WHITE)
streams = [
    ('HARDWARE', '$388–$399', 'PrecisionBrew™\nunit sale.\nTarget GM: 35–45%'),
    ('PROTOCOL PACKS', '$20–35/mo', 'Cartridge subscription.\nTarget GM: 60%+\n(Peloton benchmark: 69.1%)'),
    ('B2B / SPACES', 'Custom', 'Hotels, gyms,\ncorporate wellness.\nRevenue from Month 18+'),
    ('PLATFORM', 'TBD', 'Data insights tier,\nresearch partnerships,\nBreziIntel premium'),
]
for i, (stream, price, desc) in enumerate(streams):
    x = 0.5 + i * 3.85
    rect = s.shapes.add_shape(1, Inches(x), Inches(3.5), Inches(3.6), Inches(4.8))
    rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(18, 18, 18)
    rect.line.color.rgb = TEAL if i < 2 else RGBColor(60, 60, 60)
    tb(s, stream, x+0.15, 3.65, 3.3, 0.5, size=11, bold=True, color=TEAL)
    tb(s, price,  x+0.15, 4.2,  3.3, 0.9, size=26, bold=True, color=WHITE)
    tb(s, desc,   x+0.15, 5.2,  3.3, 2.7, size=14, color=MONO, wrap=True)

# ── S10: UNIT ECONOMICS ──────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
eyebrow(s, 'UNIT ECONOMICS')
title_line(s, 'The Subscription Is the Business.', size=44, color=WHITE)
fig = chart_unit_economics()
img_slide(s, fig, 0.5, 2.0, 7.5, 4.5)
tb(s, 'BENCHMARKS', 8.5, 2.0, 7, 0.5, size=10, color=TEAL)
benchmarks = [
    ('Keurig pods', '80.5% of US Coffee net sales ($3.2B of $4B)'),
    ('Peloton sub GM', '69.1% — hardware GM only ~13%'),
    ('Eight Sleep', '100% attach — subscription required yr 1'),
    ('Peloton churn', '1.6%/mo → 82% 12-month retention'),
    ('Target attach', '40% conservative | 60% strong case'),
    ('Protocol Pack', '$25/month → $100K MRR at 4,000 subs'),
]
for i, (label, val) in enumerate(benchmarks):
    accent_bar(s, 8.5, 2.65 + i*1.0, 0.7)
    tb(s, label, 8.7, 2.65 + i*1.0, 2.8, 0.45, size=11, bold=True, color=TEAL)
    tb(s, val,   8.7, 3.1  + i*1.0, 7.0, 0.45, size=12, color=WHITE, wrap=True)

# ── S11: REVENUE PROJECTIONS ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
eyebrow(s, 'REVENUE PROJECTIONS')
title_line(s, '18-Month Model — 10,000 Units', size=40, color=WHITE)
fig = chart_revenue()
img_slide(s, fig, 0.5, 2.0, 10, 4.5)
tb(s, 'BASE CASE (50% attach)', 11.0, 2.0, 4.7, 0.5, size=10, color=TEAL)
milestones = [
    ('$388–399', 'Hardware ASP'),
    ('$25/mo',   'Protocol Pack price'),
    ('50%',      'Subscription attach rate'),
    ('$5.4M',    'Total 18-month revenue'),
    ('$99K',     'MRR at Month 18'),
    ('1.6%/mo',  'Target monthly churn'),
    ('$399',     'LTV per sub @ 3yr tenure'),
]
for i, (num, label) in enumerate(milestones):
    tb(s, num,   11.0, 2.65 + i*0.88, 2.0, 0.5, size=18, bold=True, color=TEAL)
    tb(s, label, 13.1, 2.65 + i*0.88, 2.6, 0.5, size=12, color=MONO)

# ── S12: GO-TO-MARKET ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, GREY)
section_number(s, '03')
eyebrow(s, 'GO-TO-MARKET')
title_line(s, 'Seed 500 Right People.\nLet Them Pull the Brand Forward.', size=38, color=BLACK)
phases = [
    ('PHASE 1\nDays 1–30\n~$40K HKD',
     '• University research partnership announced\n• 15–20 Protocol Practitioners seeded (coaches, sports med, Hyrox)\n• "The Cold Protocol" newsletter launched\n• Waitlist page + 90-sec hero film'),
    ('PHASE 2\nDays 31–60\n~$120K HKD',
     '• Invite-only pop-up, Sai Ying Pun HK (200 spots)\n• Cold exposure station + live PrecisionBrew demo\n• Hero film "Start Cold" debut\n• Podcast ad buy — Huberman-adjacent APAC shows'),
    ('PHASE 3\nDays 61–90\n~$80K HKD',
     '• 30-day Baseline Protocol — structured cold+brew program\n• First 500 units at $388 USD\n• Protocol Pack subscription as post-purchase upsell\n• Referral mechanic → turn 500 into distribution network'),
]
for i, (phase, bullets) in enumerate(phases):
    x = 0.5 + i * 5.1
    rect = s.shapes.add_shape(1, Inches(x), Inches(2.2), Inches(4.9), Inches(6.3))
    c = RGBColor(43, 189, 186) if i == 0 else (RGBColor(26, 138, 136) if i == 1 else BLACK)
    rect.fill.solid(); rect.fill.fore_color.rgb = c
    rect.line.fill.background()
    tc = BLACK if i == 0 else WHITE
    tb(s, phase,   x+0.15, 2.35, 4.6, 1.1, size=13, bold=True, color=tc)
    tb(s, bullets, x+0.15, 3.6,  4.6, 4.5, size=13, color=tc, wrap=True)
tb(s, 'Total budget: ~$240K HKD (~$30K USD)   |   KPIs: 2,000 waitlist · 500 paying customers · 40% Protocol Pack attach',
   0.5, 8.5, 15, 0.5, size=11, color=MONO)

# ── S13: PRODUCT ROADMAP ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
section_number(s, '03')
eyebrow(s, 'PRODUCT ROADMAP', y=0.6)
title_line(s, '0 → 36 Months: Device → Platform → OS', size=40, color=WHITE)
fig = chart_roadmap()
img_slide(s, fig, 0.3, 2.5, 15.4, 6.0)

# ── S14: WEARABLE INTEGRATIONS ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
eyebrow(s, 'WEARABLE INTEGRATIONS')
title_line(s, 'Every Device You Already Trust.\nNow Connected to Your Morning.', size=40, color=WHITE)
fig = chart_wearable_phases()
img_slide(s, fig, 0.5, 2.5, 15, 4.0)
tb(s, 'Cronometer × WHOOP integration drove 300% increase in user engagement\nStrava × Apple Watch integration drove 20% YoY upload increase',
   0.5, 7.0, 15, 1.0, size=12, color=MONO)

# ── S15: DATA MOAT ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
eyebrow(s, 'THE DATA MOAT')
title_line(s, 'The Dataset That Cannot\nBe Reverse-Engineered.', size=44, color=WHITE)
tb(s, '3.65M', 0.6, 3.0, 8, 2.2, size=100, bold=True, color=TEAL)
tb(s, 'labeled biometric-brew-outcome events per year at 10,000 users', 0.6, 5.3, 8, 0.7, size=16, color=MONO)
moat = [
    ('Unique triple', 'Biometric state × brew parameters × 2h performance outcome — nobody else has all three'),
    ('NFC cartridge lock', 'Clones can copy the machine. They cannot copy the data layer.'),
    ('On-device ML', 'Your health data never leaves your device without consent. Privacy-first by architecture.'),
    ('University research', 'HKUST/HKU partnership → peer-reviewed publications → scientific credibility moat'),
    ('Patent portfolio', 'Two filings: thermal profiling algorithm + biometric-to-thermal mapping method'),
]
for i, (label, val) in enumerate(moat):
    accent_bar(s, 9.2, 2.1 + i*1.3, 1.1)
    tb(s, label, 9.4, 2.1  + i*1.3, 6.3, 0.5, size=13, bold=True, color=TEAL)
    tb(s, val,   9.4, 2.65 + i*1.3, 6.3, 0.6, size=13, color=WHITE, wrap=True)

# ── S16: WHY US ──────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, GREY)
section_number(s, '04')
eyebrow(s, 'WHY US')
title_line(s, 'We Are the Customer.\nThat Is the Unfair Advantage.', size=40, color=BLACK)
left = [
    ('🎓', 'Jackson (Founder) — UC Berkeley. Product-first thinking.\nThe original Protocol Optimizer.'),
    ('🏭', 'Engineering team from semiconductor + consumer electronics\n(DJI/Sony tier) — precision hardware is our DNA.'),
    ('🏃', 'Every team member is an active sports and wellness practitioner.\nWe test on ourselves before shipping to anyone.'),
    ('🔬', 'University research partnership (HKUST/HKU Medicine)\nin active discussions — credibility from day one.'),
]
right = [
    ('📍', 'Hong Kong base: Asia\'s most concentrated\nhigh-performance urban cohort. Our test market is our backyard.'),
    ('🤝', 'Graves Ho (Co-Founder): deep HK market relationships,\nmanufacturing network, operations.'),
    ('💡', 'Semiconductor background = advantage in miniaturized\nthermoelectric design most consumer brands can\'t match.'),
    ('📊', 'Phase 1 research done: 7 interviews + 181 backer surveys.\nWe know our customer. We have their data.'),
]
bullet_block(s, left,  0.5, 2.2, 7.5, 6.5, size=15, color=BLACK, gap=1.1)
bullet_block(s, right, 8.3, 2.2, 7.5, 6.5, size=15, color=BLACK, gap=1.1)

# ── S17: TRACTION ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
eyebrow(s, 'TRACTION & VALIDATION')
title_line(s, 'Not a Concept.\nA Validated Product with Early Signal.', size=40, color=WHITE)
traction = [
    ('✅', 'Kickstarter campaign successfully funded — proof of market demand'),
    ('✅', 'Live Shopify store (brezicoffee.com) + Taobao presence in China market'),
    ('✅', 'Phase 1 customer research complete: 7 in-depth interviews + 181 backer surveys'),
    ('✅', 'iOS technical spec complete. UX copy brief complete. App in development.'),
    ('✅', 'Roaster partnerships live: Regent Coffee + Stitch Coffee on brezicoffee.com'),
    ('✅', 'Hardware: Thermoelectric cooling IP validated. PrecisionBrew™ tech proven.'),
    ('✅', 'Team assembled: founder + semiconductor engineers + operations'),
    ('🔄', 'University research partnership: in active discussions (HKUST / HKU Medicine)'),
]
bullet_block(s, traction, 0.6, 2.1, 14.8, 7.0, size=17, color=WHITE, gap=0.85)

# ── S18: THE ASK ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
section_number(s, '04')
eyebrow(s, 'THE ASK', y=0.6)
title_line(s, 'Join the Protocol.', size=52, color=WHITE)
# Capital
rect1 = s.shapes.add_shape(1, Inches(0.5), Inches(2.8), Inches(7.2), Inches(5.7))
rect1.fill.solid(); rect1.fill.fore_color.rgb = RGBColor(20, 20, 20)
rect1.line.color.rgb = TEAL
tb(s, 'CAPITAL', 0.7, 2.95, 6.8, 0.5, size=11, bold=True, color=TEAL)
tb(s, '[RAISE AMOUNT]', 0.7, 3.5, 6.8, 0.9, size=32, bold=True, color=WHITE)
tb(s, 'at [VALUATION] pre-money', 0.7, 4.5, 6.8, 0.5, size=16, color=MONO)
cap_items = ['Complete precision production batch (10,000 units)',
             'iOS app + platform dev (3-month sprint)',
             'Cold Protocol Campaign — HK/SG launch',
             'University research partnership + IP filings',
             'BreziRest Pod R&D initiation']
for i, item in enumerate(cap_items):
    tb(s, f'· {item}', 0.7, 5.2 + i*0.45, 6.8, 0.4, size=13, color=WHITE)
# Talent
rect2 = s.shapes.add_shape(1, Inches(8.3), Inches(2.8), Inches(7.2), Inches(5.7))
rect2.fill.solid(); rect2.fill.fore_color.rgb = RGBColor(20, 20, 20)
rect2.line.color.rgb = CYAN2
tb(s, 'TALENT', 8.5, 2.95, 6.8, 0.5, size=11, bold=True, color=TEAL)
tb(s, 'Roles Open Now', 8.5, 3.5, 6.8, 0.7, size=26, bold=True, color=WHITE)
talent = ['Senior thermal systems hardware engineer',
          'iOS engineer (React Native + HealthKit)',
          'Product designer — industrial + UX',
          'Head of community (HK/SG Protocol network)',
          'Research coordinator — university partnership']
for i, role in enumerate(talent):
    tb(s, f'· {role}', 8.5, 4.3 + i*0.6, 6.8, 0.55, size=14, color=WHITE)

# ── S19: USE OF FUNDS ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, GREY)
eyebrow(s, 'USE OF FUNDS')
title_line(s, 'Every Dollar Has a Job.', size=44, color=BLACK)
fig = chart_use_of_funds()
img_slide(s, fig, 0.5, 2.0, 7.5, 6.5)
tb(s, 'ALLOCATION RATIONALE', 8.5, 2.0, 7.2, 0.5, size=10, color=TEAL)
allocs = [
    ('35%  Hardware Production',   'First 10,000-unit precision batch. NFC cartridge tooling.'),
    ('25%  App & Platform Dev',     '3-month iOS sprint + backend data infrastructure + Apple Health integration.'),
    ('20%  GTM & Brand',            '90-day Cold Protocol Campaign. Pop-up. Hero film. Podcast media.'),
    ('12%  R&D & Patents',          'Two patent filings. BreziRest Pod early R&D. University research seed funding.'),
    (' 8%  Operations & Team',      'Hiring critical roles. Legal. Corporate structure in HK.'),
]
for i, (alloc, detail) in enumerate(allocs):
    accent_bar(s, 8.5, 2.65 + i*1.2, 1.0)
    tb(s, alloc,  8.7, 2.65 + i*1.2, 7.0, 0.45, size=13, bold=True, color=BLACK)
    tb(s, detail, 8.7, 3.15 + i*1.2, 7.0, 0.55, size=13, color=MONO, wrap=True)

# ── S20: CLOSING ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, BLACK)
tb(s, 'bre:zi', 2, 2.5, 12, 2.0, size=96, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 'Brewed for your biology.', 2, 4.8, 12, 0.9, size=28, color=TEAL, align=PP_ALIGN.CENTER)
tb(s, 'brezicoffee.com', 2, 6.0, 12, 0.5, size=16, color=MONO, align=PP_ALIGN.CENTER)
tb(s, 'jackson@brezi.com  ·  Hong Kong', 2, 6.6, 12, 0.5, size=14, color=MONO, align=PP_ALIGN.CENTER)

# ── SAVE ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f'✅ Saved: {OUT}  ({prs.slides.__len__()} slides)')

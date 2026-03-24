#!/usr/bin/env python3
"""Build Brezi Cold Protocol Investor/Talent Deck v2 — 18 slides, 16x9."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ── Colours ──────────────────────────────────────────────────────────────
BLACK       = RGBColor(0x00, 0x00, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE     = RGBColor(0xF5, 0xF5, 0xF5)
BONE        = RGBColor(0xF0, 0xEB, 0xE1)
TEAL        = RGBColor(0x2B, 0xBD, 0xBA)
CRYO        = RGBColor(0x1A, 0x8A, 0x88)
MONUMENT    = RGBColor(0x8C, 0x8C, 0x8C)

FONT_SANS = "Helvetica Neue"

# ── Helpers ──────────────────────────────────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *,
             font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
             font_name=FONT_SANS, line_spacing=None, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox


def add_bullets(slide, left, top, width, height, items, *,
                font_size=16, color=BLACK, bullet_char="•",
                line_spacing=24, font_name=FONT_SANS):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char}  {item}" if bullet_char else item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
    return txBox


def add_accent_bar(slide, left, top, height, width=0.06, color=TEAL):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_eyebrow(slide, left, top, text, color=TEAL, width=6):
    return add_text(slide, left, top, width, 0.4, text,
                    font_size=10, bold=True, color=color)


def add_stat(slide, left, top, width, height, number, label, *,
             num_size=80, num_color=WHITE, label_color=MONUMENT,
             label_size=16, alignment=PP_ALIGN.LEFT):
    add_text(slide, left, top, width, height * 0.6, number,
             font_size=num_size, bold=True, color=num_color, alignment=alignment)
    add_text(slide, left, top + height * 0.55, width, height * 0.4, label,
             font_size=label_size, color=label_color, alignment=alignment)


def add_multirun_text(slide, left, top, width, height, runs, *,
                      alignment=PP_ALIGN.LEFT, line_spacing=None):
    """runs = list of (text, size, bold, color)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    for text, size, bold, color in runs:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT_SANS
    return txBox


# ── Build deck ───────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
blank_layout = prs.slide_layouts[6]  # blank


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, 0.8, 0.5, 6, 0.4, "SAI YING PUN · HONG KONG · 2026",
         font_size=10, color=MONUMENT)
add_text(s, 2, 2.8, 12, 2, "bre:zi",
         font_size=96, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(s, 2, 4.8, 12, 1, "Brewed for your biology.",
         font_size=28, color=TEAL, alignment=PP_ALIGN.CENTER)
add_text(s, 2, 6.5, 12, 0.6, "Investor & Talent Deck — Confidential",
         font_size=14, color=MONUMENT, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — SECTION HEADER: "The World Has Changed"
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, 0.8, 0.6, 2, 0.6, "01", font_size=32, bold=True, color=TEAL)
add_text(s, 1.5, 3, 13, 2, "The World Has Changed",
         font_size=64, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — CONTENT DARK: "The Protocol Optimizer Exists at Scale"
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_eyebrow(s, 0.8, 0.6, "THE OPPORTUNITY")
add_text(s, 0.8, 1.2, 14, 1, "The Protocol Optimizer Exists at Scale",
         font_size=44, bold=True, color=WHITE)

add_stat(s, 0.8, 2.8, 6, 1.8, "538M", "wearable units shipped in 2024",
         num_size=72, num_color=TEAL)
add_stat(s, 0.8, 4.8, 6, 1.5, "32%", "of wearable owners have attached subscriptions",
         num_size=60, num_color=WHITE)
add_stat(s, 8, 4.8, 6, 1.5, "92%", "would pay more for at least one health feature",
         num_size=60, num_color=WHITE)

add_text(s, 0.8, 7.8, 14, 0.5,
         "Source: IDC 2024 / Parks Associates consumer research",
         font_size=10, color=MONUMENT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — CONTENT LIGHT: "Meet the Protocol Optimizer"
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, SURFACE)
add_eyebrow(s, 0.8, 0.6, "TARGET CUSTOMER")
add_text(s, 0.8, 1.2, 10, 1, "Meet the Protocol Optimizer",
         font_size=44, bold=True, color=BLACK)
add_accent_bar(s, 0.8, 2.4, 4.5)

bullets = [
    "Age 27–38 | Urban | HHI $200K+ | Advanced degree",
    "Owns Oura Ring or WHOOP. Trains Hyrox, Zone 2, or CrossFit 4×/week",
    "Follows Huberman Lab, Peter Attia. Has read Outlive.",
    "Spends $400 on the right thing without hesitation",
    "Healthfulness overtakes price as #1 purchase driver at $100K+ income",
    "Discovery: Podcast → peer community → purchase",
]
add_bullets(s, 1.1, 2.4, 9, 5, bullets, color=BLACK)

# Right-side callout
add_accent_bar(s, 11, 3, 3, width=0.08)
add_text(s, 11.3, 3.2, 4, 2.5,
         'At ≥$100K income, healthfulness overtakes price as #1 purchase driver\n(IFIC 2024)',
         font_size=14, color=MONUMENT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — CONTENT LIGHT: "The Gap"
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, SURFACE)
add_eyebrow(s, 0.8, 0.6, "THE PROBLEM")
add_text(s, 0.8, 1.2, 14, 1.2,
         "Everything Is Optimized.\nExcept the First Thing They Consume.",
         font_size=44, bold=True, color=BLACK)
add_accent_bar(s, 0.8, 2.8, 5)

bullets = [
    "Morning stack: sleep tracker ✓, training plan ✓, cold exposure ✓, supplements ✓, caffeine timing... ✗",
    "Cold brew is the #1 anchor ritual of this cohort — yet made with tools designed for convenience, not protocol",
    "No product closes the loop between cold exposure + morning consumption + biometric feedback",
    "The ritual is fragmented across 4–6 unconnected products",
    "Cold plunge market: $330M (2024) → $660M by 2033 — and nobody owns the adjacent morning ritual",
]
add_bullets(s, 1.1, 2.8, 13, 5, bullets, color=BLACK, line_spacing=28)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — QUOTE
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BONE)
add_text(s, 2, 2.5, 12, 3,
         '"Your best days don\'t happen by accident."',
         font_size=40, bold=False, color=BLACK, alignment=PP_ALIGN.CENTER)
add_text(s, 2, 5.5, 12, 1,
         "— The belief that drives every Protocol Optimizer.\n   And every Brezi design decision.",
         font_size=16, color=MONUMENT, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — SECTION HEADER: "Introducing bre:zi"
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, 0.8, 0.6, 2, 0.6, "02", font_size=32, bold=True, color=TEAL)
add_text(s, 1.5, 3, 13, 2, "Introducing bre:zi",
         font_size=64, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — CONTENT DARK: "The Cold Protocol Platform"
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_eyebrow(s, 0.8, 0.6, "THE PRODUCT")
add_text(s, 0.8, 1.2, 14, 1,
         "bre:zi — The Cold Protocol Platform",
         font_size=44, bold=True, color=WHITE)
add_accent_bar(s, 0.8, 2.6, 4.5)

bullets = [
    "PrecisionBrew™: thermoelectric cold extraction in 12 minutes — no ice, no heat, 67% less acidic",
    "Reads your biometric state before every brew. Adapts the protocol to your body today.",
    "NFC cartridge system: every Protocol Pack is formula-validated and data-encoded",
    "Hardware + App + Platform: the first cold wellness instrument that learns your biology",
    "Price: $388–$399  |  Protocol Packs: $20–35/month",
]
add_bullets(s, 1.1, 2.6, 13, 5, bullets, color=WHITE, line_spacing=28)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — CONTENT LIGHT: "Object. Ritual. Temperature."
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, SURFACE)
add_eyebrow(s, 0.8, 0.6, "BRAND WORLD")
add_text(s, 0.8, 1.2, 14, 1, "Object. Ritual. Temperature.",
         font_size=44, bold=True, color=BLACK)

brand_items = [
    ("IT LOOKS LIKE:", "Matte anodized aluminum, frosted glass. Camera energy. Watch energy. Not blender energy."),
    ("IT FEELS LIKE:", "Protocol calm. 5:47am. Your system initializing. Cold coming. Focus coming."),
    ("IT SOUNDS LIKE:", "One soft click when the cycle locks. Then 12 minutes of silence."),
    ("IT SITS NEXT TO:", "Your WHOOP. Your Fellow Stagg. Your Oura Ring charger. It passed the shelf test."),
    ("IT IS NOT:", 'Wellness in the lavender sense. Not for people who say "self-care." For people who say "protocol."'),
]

y = 2.6
for label, desc in brand_items:
    add_accent_bar(s, 0.8, y, 0.9)
    add_text(s, 1.1, y, 3, 0.5, label, font_size=14, bold=True, color=CRYO)
    add_text(s, 4.0, y, 11, 0.9, desc, font_size=16, color=BLACK)
    y += 1.15


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — CONTENT DARK: "The 12-Minute Window"
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_eyebrow(s, 0.8, 0.6, "THE RITUAL")

# Build "12:00" with teal colon using multi-run
txBox = s.shapes.add_textbox(Inches(2), Inches(1.8), Inches(12), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
for txt, clr in [("12", WHITE), (":", TEAL), ("00", WHITE)]:
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(120)
    run.font.bold = True
    run.font.color.rgb = clr
    run.font.name = FONT_SANS

add_text(s, 1.5, 4.5, 13, 2.5,
         "The brew takes 12 minutes. That is not a limitation — that is the product. "
         "A single cold-blue indicator light makes one slow revolution. "
         "No progress bar. No notification. Just the light. Moving. Once around in 12 minutes.",
         font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER, line_spacing=28)

add_text(s, 1.5, 7.5, 13, 0.6,
         "The 12-minute window is not a feature. It is the protocol promise made physical.",
         font_size=12, color=MONUMENT, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — SECTION HEADER: "Platform & Data Moat"
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, 0.8, 0.6, 2, 0.6, "03", font_size=32, bold=True, color=TEAL)
add_text(s, 1.5, 3, 13, 2, "Platform & Data Moat",
         font_size=64, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — CONTENT LIGHT: "Brezi OS — Three Layers"
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, SURFACE)
add_eyebrow(s, 0.8, 0.6, "TECHNOLOGY")
add_text(s, 0.8, 1.2, 14, 1.2,
         "Every Brew Is a Data Event.\nEvery Event Makes the Next Brew Smarter.",
         font_size=40, bold=True, color=BLACK)

layers = [
    ("DEVICE LAYER",
     "PrecisionBrew™ captures formula data, extraction curve, timestamp — structured data, every brew"),
    ("HEALTH GRAPH",
     "Links each brew to your biometric state — HRV, sleep score, recovery. Closes the loop with your 2-hour performance log."),
    ("BREZINTEL",
     '"Your HRV is down 18% — try the Recovery Protocol today." Adaptive. Private. On-device.'),
]

x_pos = 0.8
col_w = 4.5
for label, desc in layers:
    add_accent_bar(s, x_pos, 3.2, 3.5)
    add_text(s, x_pos + 0.25, 3.3, col_w, 0.5, label,
             font_size=14, bold=True, color=CRYO)
    add_text(s, x_pos + 0.25, 3.9, col_w, 3, desc,
             font_size=16, color=BLACK, line_spacing=24)
    x_pos += col_w + 0.4


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 — CONTENT DARK: "The Data Moat"
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_eyebrow(s, 0.8, 0.6, "DEFENSIBILITY")
add_text(s, 0.8, 1.2, 14, 1,
         "The Dataset That Cannot Be Reverse-Engineered",
         font_size=44, bold=True, color=WHITE)

add_text(s, 0.8, 2.6, 7, 1.2, "3.65M",
         font_size=96, bold=True, color=TEAL)
add_text(s, 0.8, 4.0, 7, 0.6,
         "labeled biometric-brew-outcome events per year at 10K users",
         font_size=16, color=MONUMENT)

moat_bullets = [
    "Unique triple: biometric state × brew parameters × 2h subjective performance — no other platform has all three",
    "NFC cartridge authentication: clones can copy the machine, not the data layer",
    "On-device ML: your health data never leaves your device without consent — privacy-first by architecture",
    "University research partnership (HKUST/HKU): peer-reviewed publications validate the dataset scientifically",
    "Two patent filings: thermal profiling algorithm + biometric-to-thermal mapping method",
]
add_bullets(s, 0.8, 4.8, 14, 3.5, moat_bullets, color=WHITE, line_spacing=26)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 — CONTENT LIGHT: "Wearable Integration Roadmap"
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, SURFACE)
add_eyebrow(s, 0.8, 0.6, "INTEGRATIONS")
add_text(s, 0.8, 1.2, 14, 1,
         "Every Device You Trust. Now Connected to Your Morning.",
         font_size=40, bold=True, color=BLACK)

phases = [
    ("PHASE 1 (NOW)", "Apple HealthKit — 1B+ users, zero rate limits, ships Day 1"),
    ("PHASE 2 (Month 6)", "Oura Ring — best sleep/HRV data, free OAuth2 API"),
    ("PHASE 3 (Month 12)", "WHOOP — 300% engagement lift benchmark, recovery protocol users"),
    ("PHASE 4 (Month 18)", "Ultrahuman — APAC priority, CGM metabolic data, invite-only partnership"),
]

x_pos = 0.8
col_w = 3.5
for label, desc in phases:
    add_accent_bar(s, x_pos, 3, 3.5)
    add_text(s, x_pos + 0.2, 3.1, col_w, 0.6, label,
             font_size=13, bold=True, color=CRYO)
    add_text(s, x_pos + 0.2, 3.7, col_w, 3, desc,
             font_size=15, color=BLACK, line_spacing=22)
    x_pos += col_w + 0.3


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15 — SECTION HEADER: "Business Model"
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_text(s, 0.8, 0.6, 2, 0.6, "04", font_size=32, bold=True, color=TEAL)
add_text(s, 1.5, 3, 13, 2, "Business Model",
         font_size=64, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16 — CONTENT DARK: "Unit Economics"
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_eyebrow(s, 0.8, 0.6, "THE NUMBERS")
add_text(s, 0.8, 1.2, 14, 1,
         "Hardware Opens the Door. Protocol Packs Own the Room.",
         font_size=44, bold=True, color=WHITE)

# Key stats in a 2x2 grid
stats = [
    ("$388–$399", "Hardware ASP"),
    ("$20–35/mo", "Protocol Packs (Keurig: pods = 80.5% of revenue)"),
    ("60%+", "Subscription GM target (Peloton: 69.1%)"),
    ("≤1.6%", "Monthly churn target (Peloton: 82% 12-mo retention)"),
]
positions = [(0.8, 2.6), (8, 2.6), (0.8, 4.2), (8, 4.2)]
for (num, label), (x, y) in zip(stats, positions):
    add_text(s, x, y, 6, 0.7, num, font_size=48, bold=True, color=TEAL)
    add_text(s, x, y + 0.7, 6, 0.5, label, font_size=14, color=MONUMENT)

# 18-month model table header
table_top = 5.8
add_text(s, 0.8, table_top - 0.5, 14, 0.5, "18-MONTH REVENUE MODEL",
         font_size=12, bold=True, color=TEAL)

# Build table
tbl = s.shapes.add_table(4, 5, Inches(0.8), Inches(table_top),
                          Inches(14), Inches(2)).table
tbl.columns[0].width = Inches(2.8)
tbl.columns[1].width = Inches(2.8)
tbl.columns[2].width = Inches(2.8)
tbl.columns[3].width = Inches(2.8)
tbl.columns[4].width = Inches(2.8)

headers = ["Scenario", "Units", "Attach Rate", "MRR (Mo 18)", "Total Revenue"]
rows_data = [
    ["Conservative", "10,000", "40%", "$79.5K", "$5.1M"],
    ["Base", "10,000", "50%", "$99.4K", "$5.6M"],
    ["Strong", "10,000", "60%", "$119.2K", "$5.65M"],
]

for col_i, h in enumerate(headers):
    cell = tbl.cell(0, col_i)
    cell.text = h
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.font.name = FONT_SANS
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

for row_i, row in enumerate(rows_data):
    for col_i, val in enumerate(row):
        cell = tbl.cell(row_i + 1, col_i)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.font.name = FONT_SANS
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x0D)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 17 — CONTENT LIGHT: "Team & Advisors"
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, SURFACE)
add_eyebrow(s, 0.8, 0.6, "THE TEAM")
add_text(s, 0.8, 1.2, 14, 1, "Built by the People Who Use It",
         font_size=44, bold=True, color=BLACK)
add_accent_bar(s, 0.8, 2.6, 5)

team_bullets = [
    "Jackson (Founder) — UC Berkeley | Product vision | Protocol Optimizer himself",
    "Graves Ho (Co-Founder) — Operations | HK market | Manufacturing relationships",
    "Engineering Team — Semiconductor + consumer electronics background (DJI / Sony tier)",
    "All team members: active sports and wellness practitioners. We are the customer.",
    "Advisory: University research partner (HKUST/HKU Medicine) — in discussions",
    "Advisory: Industrial design and thermal engineering expertise — assembled",
]
add_bullets(s, 1.1, 2.6, 13, 5, team_bullets, color=BLACK, line_spacing=28)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 18 — CONTENT DARK: "The Ask"
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BLACK)
add_eyebrow(s, 0.8, 0.6, "JOIN THE PROTOCOL")
add_text(s, 0.8, 1.2, 14, 1, "What We're Building Together",
         font_size=44, bold=True, color=WHITE)

# Left column — Talent
add_text(s, 0.8, 2.6, 6, 0.5, "TALENT WE NEED",
         font_size=14, bold=True, color=TEAL)
add_accent_bar(s, 0.8, 3.2, 3.5)
talent = [
    "Senior thermal systems hardware engineer",
    "iOS engineer (React Native, health APIs)",
    "Product designer — industrial + UX",
    "Performance community builder (HK/SG)",
]
add_bullets(s, 1.1, 3.2, 6, 3.5, talent, color=WHITE, line_spacing=28)

# Right column — Capital
add_text(s, 8.5, 2.6, 6, 0.5, "CAPITAL WE'RE RAISING",
         font_size=14, bold=True, color=TEAL)
add_accent_bar(s, 8.5, 3.2, 3.5)
capital = [
    "18-month runway: precision batch + app build + BreziRest R&D",
    "GTM: Cold Protocol Campaign ($30K USD seed round)",
    "Research partnership + IP filings",
    "Target: [RAISE AMOUNT] at [VALUATION]",
]
add_bullets(s, 8.8, 3.2, 6, 3.5, capital, color=WHITE, line_spacing=28)

# Footer
add_text(s, 2, 7.8, 12, 0.6,
         "brezicoffee.com  |  Jackson Hung  |  jackson@brezi.com",
         font_size=16, color=MONUMENT, alignment=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "Brezi_Cold_Protocol_Deck_v2.pptx")
prs.save(out_path)
print(f"✓ Saved {out_path}  ({len(prs.slides)} slides)")
